# -*- coding: utf-8 -*-
"""Generate the iFinance Implementation Plan executive deck (PPTX).

Mirrors the HTML plan artifact v1.2: 8 modules, agile 2-week sprints,
Service Pricing in Sprint 2 alongside Reporting & BI.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------- palette (validated in the dataviz pass) ----------
BRAND    = RGBColor.from_string("0F5262")   # deep teal
PLATFORM = RGBColor.from_string("0E7EA3")   # gantt platform bars
MODULE   = RGBColor.from_string("B0721E")   # gantt module bars
ACCENT   = RGBColor.from_string("A06612")   # gold accents
INK      = RGBColor.from_string("1C2B31")
MUTED    = RGBColor.from_string("55676E")
FAINT    = RGBColor.from_string("7E9096")
LINE     = RGBColor.from_string("DCE4E6")
SOFT     = RGBColor.from_string("E3EEF1")   # brand soft fill
ACC_SOFT = RGBColor.from_string("F4EAD8")
NEUTRAL  = RGBColor.from_string("9AA9AF")   # de-emphasised bars
TRACK    = RGBColor.from_string("EDF1F2")
WHITE    = RGBColor.from_string("FFFFFF")
GROUND   = RGBColor.from_string("F8FAFA")

SERIF = "Georgia"
SANS  = "Segoe UI"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide_new(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, round_=False, radius=0.06):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def tb(s, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf

def para(tf, runs, size=12, color=INK, bold=False, font=SANS, align=PP_ALIGN.LEFT,
         space_after=4, space_before=0, first=False, line_spacing=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line_spacing: p.line_spacing = line_spacing
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        r = p.add_run(); r.text = text
        r.font.name = ov.get("font", font)
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.color.rgb = ov.get("color", color)
        if ov.get("italic"): r.font.italic = True
    return p

def bullets(tf, items, size=13, gap=8, first=True, lead_color=INK, body_color=MUTED):
    """items: list of (lead, rest) or plain string."""
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            lead, rest = it
            runs = [("•  ", {"color": ACCENT, "bold": True}),
                    (lead, {"bold": True, "color": lead_color}),
                    (rest, {"color": body_color})]
        else:
            runs = [("•  ", {"color": ACCENT, "bold": True}),
                    (it, {"color": body_color})]
        para(tf, runs, size=size, space_after=gap, first=(first and i == 0), line_spacing=1.15)

def header(s, sec_no, title, sub=None):
    tf = tb(s, 0.55, 0.32, 12.2, 0.9)
    para(tf, [(f"{sec_no}   ", {"font": SERIF, "size": 20, "bold": True, "color": ACCENT}),
              (title, {"font": SERIF, "size": 26, "bold": True, "color": INK})], first=True)
    rect(s, 0.55, 1.02, 12.23, 0.02, fill=LINE)
    if sub:
        tf2 = tb(s, 0.55, 1.12, 12.2, 0.35)
        para(tf2, sub, size=12, color=MUTED, first=True)

def footer(s, n):
    tf = tb(s, 0.55, 7.08, 9.0, 0.3)
    para(tf, "iFinance Implementation Plan  ·  DCT Finance Department  ·  July 2026",
         size=8.5, color=FAINT, first=True)
    tf2 = tb(s, 12.3, 7.08, 0.5, 0.3)
    para(tf2, str(n), size=8.5, color=FAINT, align=PP_ALIGN.RIGHT, first=True)

# =========================================================
# Slide 1 — Title
# =========================================================
s = slide_new(bg=BRAND)
rect(s, 0, 0, 13.333, 0.12, fill=ACCENT)
tf = tb(s, 1.0, 1.7, 11.3, 0.5)
para(tf, "DEPARTMENT OF CULTURE AND TOURISM — ABU DHABI   ·   FINANCE DEPARTMENT",
     size=13, color=RGBColor.from_string("D9A441"), bold=True, first=True)
tf = tb(s, 1.0, 2.25, 11.3, 1.4)
para(tf, "iFinance Implementation Plan", font=SERIF, size=48, bold=True, color=WHITE, first=True)
tf = tb(s, 1.0, 3.6, 10.6, 1.2)
para(tf, "Deployment on a dedicated Oracle OCI tenancy within the Core42 data centre — "
         "co-located with the ADERP Oracle Fusion instance and fully integrated with it in both directions.",
     size=16, color=RGBColor.from_string("C9DEE4"), first=True, line_spacing=1.2)
rect(s, 1.0, 5.15, 3.2, 0.025, fill=RGBColor.from_string("D9A441"))
tf = tb(s, 1.0, 5.45, 11.3, 1.4)
para(tf, [("Submitted to:  ", {"bold": True, "color": WHITE}),
          ("Finance Director — for review & approval", {"color": RGBColor.from_string("C9DEE4")})],
     size=13, first=True, space_after=6)
para(tf, [("Methodology:  ", {"bold": True, "color": WHITE}),
          ("Agile delivery — two-week sprints", {"color": RGBColor.from_string("C9DEE4")})],
     size=13, space_after=6)
para(tf, [("Document:  ", {"bold": True, "color": WHITE}),
          ("Implementation Plan v1.2  ·  8 July 2026  ·  Prepared by the Finance Systems Team",
           {"color": RGBColor.from_string("C9DEE4")})], size=13)

# =========================================================
# Slide 2 — Executive summary
# =========================================================
s = slide_new()
header(s, "", "Executive Summary — Recommendation")
rect(s, 0.55, 1.35, 0.06, 5.35, fill=BRAND)
tf = tb(s, 0.95, 1.5, 11.7, 5.3)
bullets(tf, [
    ("Solution:  ", "the iFinance platform, deployed on a dedicated Oracle OCI tenancy provisioned within the "
     "Core42 data centre — the same facility hosting the ADERP Oracle Fusion instance."),
    ("Investment:  ", "a fixed quotation of USD 100,000, received for the entire eight-module scope — "
     "General Ledger, Reporting, Service Pricing, Freelancers, Team Management, Petty Cash, Duty Travel "
     "and Employee Housing."),
    ("Methodology:  ", "Agile delivery — two-week sprints, each concluding with user acceptance testing (UAT), "
     "training and go-live; Sprint 2 delivers Reporting & BI and Service Pricing together."),
    ("Timeline:  ", "tenancy provisioned within one week of contract signature; the first module (General Ledger) "
     "live in week 3; a further module every two weeks; full scope live in ≈ 3 months."),
    ("Rationale:  ", "the technology best aligned with ADGE and Oracle Fusion  ·  full bidirectional Fusion "
     "integration  ·  the fastest achievable delivery  ·  the lowest total cost of the available options."),
    ("Critical path:  ", "a single enabling action — procurement of the Oracle OCI tenancy within Core42."),
], size=15, gap=14)
footer(s, 2)

# =========================================================
# Slide 3 — Key figures
# =========================================================
s = slide_new()
header(s, "", "Key Figures")
kpis = [
    ("USD 100K", "Fixed quotation — entire 8-module scope"),
    ("8", "Business modules delivered"),
    ("≈ 3 mo", "Full scope live, from contract signature"),
    ("Week 3", "First module live — General Ledger"),
    ("~12.5K", "Average cost per module (USD)"),
    ("100%", "Data resident in Abu Dhabi (Core42)"),
]
kx, ky, kw, kh, gx, gy = 0.55, 1.6, 3.94, 2.15, 0.2, 0.35
for i, (v, l) in enumerate(kpis):
    col, row = i % 3, i // 3
    x = kx + col * (kw + gx); y = ky + row * (kh + gy)
    rect(s, x, y, kw, kh, fill=GROUND, line=LINE, round_=True, radius=0.08)
    rect(s, x, y, 0.05, kh, fill=ACCENT)
    tf = tb(s, x + 0.3, y + 0.42, kw - 0.55, 0.85)
    para(tf, v, font=SERIF, size=34, bold=True, color=BRAND, first=True)
    tf = tb(s, x + 0.3, y + 1.3, kw - 0.55, 0.7)
    para(tf, l, size=12.5, color=MUTED, first=True, line_spacing=1.1)
footer(s, 3)

# =========================================================
# Slide 4 — Why this option (4 pillars)
# =========================================================
s = slide_new()
header(s, "1", "Why This Option")
pillars = [
    ("TECHNOLOGY ALIGNMENT", "Best technology — aligned with ADGE & Oracle Fusion", [
        "End-to-end Oracle stack — OCI, Autonomous Database, APEX/ORDS, Oracle JET.",
        "Same technology family, and the same Core42 data centre, as ADERP Fusion.",
        "Conforms to ADGE cloud & data-residency direction — no new vendor.",
    ]),
    ("FUSION INTEGRATION — BOTH DIRECTIONS", "Fully integrated with Oracle Fusion", [
        "Outbound: scheduled extracts of GL, budgets, projects, AP / PO / PR — already operational in pilot.",
        "Inbound: approved transactions (e.g. petty-cash reimbursements) post back to Fusion AP.",
        "Co-location keeps integration traffic inside the data centre — data never leaves Abu Dhabi.",
    ]),
    ("SPEED OF DELIVERY", "Very fast delivery, at the pace the business requires", [
        "Agile — two-week sprints, inclusive of UAT and training.",
        "Six of eight modules already built and UAT-tested in the pilot environment.",
        "First go-live in week 3; a further module every two weeks.",
    ]),
    ("COST", "Lowest total cost", [
        "USD 100,000 all-inclusive (~USD 12.5K per module) — deployment to hypercare.",
        "No per-user licensing — APEX / ORDS / JET runtimes licence-included.",
        "Market alternatives indicatively 5–10× the cost and 4–6× the duration.",
    ]),
]
px, py, pw, ph, pgx, pgy = 0.55, 1.35, 6.0, 2.7, 0.23, 0.22
for i, (tag, title, items) in enumerate(pillars):
    col, row = i % 2, i // 2
    x = px + col * (pw + pgx); y = py + row * (ph + pgy)
    rect(s, x, y, pw, ph, fill=WHITE, line=LINE, round_=True, radius=0.05)
    tf = tb(s, x + 0.28, y + 0.2, pw - 0.55, 0.3)
    para(tf, tag, size=9.5, bold=True, color=ACCENT, first=True)
    tf = tb(s, x + 0.28, y + 0.5, pw - 0.55, 0.5)
    para(tf, title, font=SERIF, size=15.5, bold=True, color=BRAND, first=True)
    tf = tb(s, x + 0.28, y + 1.02, pw - 0.55, ph - 1.15)
    bullets(tf, items, size=10.5, gap=4)
footer(s, 4)

# =========================================================
# Slide 5 — Solution architecture
# =========================================================
s = slide_new()
header(s, "2", "Solution Architecture")
# user chips
chips = ["Finance staff — web (EN / AR)", "Approvers — mobile app + push",
         "Management — dashboards & BI", "Freelancers — public portal"]
cw = 2.85; cx = (13.333 - (cw * 4 + 0.25 * 3)) / 2
for i, c in enumerate(chips):
    x = cx + i * (cw + 0.25)
    rect(s, x, 1.32, cw, 0.42, fill=SOFT, round_=True, radius=0.5)
    tf = tb(s, x, 1.32, cw, 0.42, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, c, size=10.5, bold=True, color=BRAND, align=PP_ALIGN.CENTER, first=True, space_after=0)
tf = tb(s, 0.55, 1.83, 12.23, 0.3)
para(tf, "Single sign-on across all interfaces  ·  role-based access  ·  full audit trail",
     size=9.5, color=FAINT, align=PP_ALIGN.CENTER, first=True)
# connector
con = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.53), Inches(2.12), Inches(0.28), Inches(0.32))
con.fill.solid(); con.fill.fore_color.rgb = FAINT; con.line.fill.background(); con.shadow.inherit = False
# DC box
rect(s, 0.85, 2.55, 11.63, 3.7, fill=None, line=BRAND, line_w=2.0, round_=True, radius=0.035)
tf = tb(s, 1.15, 2.72, 11.0, 0.32)
para(tf, "CORE42 DATA CENTRE — ABU DHABI", size=11, bold=True, color=BRAND, first=True)
# left sysbox
rect(s, 1.15, 3.12, 4.6, 2.85, fill=GROUND, line=LINE, round_=True, radius=0.05)
tf = tb(s, 1.4, 3.3, 4.15, 0.45)
para(tf, "DCT iFinance — dedicated OCI tenancy", font=SERIF, size=13.5, bold=True, color=INK, first=True)
tf = tb(s, 1.4, 3.78, 4.15, 2.1)
bullets(tf, ["Oracle Autonomous Database (data + business logic)",
             "APEX & ORDS REST services",
             "Web tier (application delivery)",
             "Integration runners — extract & write-back",
             "Reporting engines — scheduled PDF / Excel + interactive BI"], size=10, gap=4)
# right sysbox
rect(s, 7.75, 3.12, 4.6, 2.85, fill=GROUND, line=LINE, round_=True, radius=0.05)
tf = tb(s, 8.0, 3.3, 4.15, 0.45)
para(tf, "ADERP — Oracle Fusion Cloud", font=SERIF, size=13.5, bold=True, color=INK, first=True)
tf = tb(s, 8.0, 3.78, 4.15, 2.1)
bullets(tf, ["Government shared ERP instance",
             "Financials — GL, AP, PO / procurement",
             "Projects & budgets",
             "OTBI / BIP reporting layer (extract source)"], size=10, gap=4)
# arrows between
a1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.95), Inches(3.7), Inches(1.6), Inches(0.3))
a1.fill.solid(); a1.fill.fore_color.rgb = BRAND; a1.line.fill.background(); a1.shadow.inherit = False
tf = tb(s, 5.7, 4.03, 2.1, 0.55)
para(tf, "Approved transactions\n(e.g. AP invoices)", size=8.5, color=MUTED, align=PP_ALIGN.CENTER, first=True)
a2 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(5.95), Inches(4.85), Inches(1.6), Inches(0.3))
a2.fill.solid(); a2.fill.fore_color.rgb = PLATFORM; a2.line.fill.background(); a2.shadow.inherit = False
tf = tb(s, 5.65, 5.18, 2.2, 0.6)
para(tf, "Scheduled data extracts\n(GL, budgets, AP / PO / PR)", size=8.5, color=MUTED, align=PP_ALIGN.CENTER, first=True)
# caption
tf = tb(s, 0.85, 6.4, 11.63, 0.4)
para(tf, "Single dedicated tenancy  ·  integration traffic remains inside the data centre  ·  "
         "all data resident in Abu Dhabi  ·  government security baseline applied at provisioning",
     size=10, color=FAINT, align=PP_ALIGN.CENTER, first=True)
footer(s, 5)

# =========================================================
# Slide 6 — Scope & investment
# =========================================================
s = slide_new()
header(s, "3", "Scope & Investment",
       sub="USD 100,000 fixed quotation — eight business modules in delivery order, on a platform foundation included at no additional cost.")
rows = [
    ("1", "General Ledger", "CoA classifications (sector / chapter / programme), budget-vs-actuals and budget-utilization reporting over ADERP data", "Built & piloted"),
    ("2", "Reporting & BI", "Scheduled report distribution (PDF / Excel by email) and interactive self-service reports", "Built & piloted"),
    ("3", "Service Pricing", "Pricing catalogue for DCT services — price definition, change approvals, effective-dated price lists (Sprint 2, with Reporting)", "New build"),
    ("4", "Freelancers", "Registration (incl. public portal with AI-assisted document capture), contracts, vouchers, payments", "Built & piloted"),
    ("5", "Team Management", "Team tasks, periodic reporting cycles, member check-ins, executive dashboards", "Built & piloted"),
    ("6", "Petty Cash", "Custody, expenses, clearing & reimbursement with approvals; posts to Fusion AP", "Built & piloted"),
    ("7", "Duty Travel", "Travel requests, per-diem calculation, approvals and settlement", "Built & piloted"),
    ("8", "Employee Housing", "Housing allowance & accommodation requests, eligibility, approvals, payment instructions", "New build"),
    ("", "Platform foundation", "Included: SSO, roles, delegations, audit, unified approvals inbox, notifications, bilingual EN / AR, mobile app", "Included"),
]
tbl_shape = s.shapes.add_table(len(rows) + 1, 4, Inches(0.55), Inches(1.62), Inches(12.23), Inches(4.7))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(0.5)
tbl.columns[1].width = Inches(2.3)
tbl.columns[2].width = Inches(7.63)
tbl.columns[3].width = Inches(1.8)
hdr = ["#", "MODULE", "WHAT IT DELIVERS", "STATUS TODAY"]
for j, htxt in enumerate(hdr):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = BRAND
    c.margin_left = Inches(0.08); c.margin_right = Inches(0.05)
    c.margin_top = c.margin_bottom = Inches(0.03)
    tf = c.text_frame; tf.word_wrap = True
    para(tf, htxt, size=10, bold=True, color=WHITE, first=True, space_after=0)
for i, (n, m, d, st) in enumerate(rows, start=1):
    vals = [n, m, d, st]
    fill = WHITE if i % 2 else GROUND
    if m == "Platform foundation": fill = SOFT
    for j, v in enumerate(vals):
        c = tbl.cell(i, j)
        c.fill.solid(); c.fill.fore_color.rgb = fill
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.08); c.margin_right = Inches(0.05)
        c.margin_top = c.margin_bottom = Inches(0.02)
        tf = c.text_frame; tf.word_wrap = True
        if j == 3:
            col = MODULE if v == "New build" else BRAND
            para(tf, v, size=9.5, bold=True, color=col, first=True, space_after=0)
        elif j == 1:
            para(tf, v, size=10, bold=True, color=INK, first=True, space_after=0)
        else:
            para(tf, v, size=9.5, color=MUTED, first=True, space_after=0)
tf = tb(s, 0.55, 6.5, 12.23, 0.45)
para(tf, "The quotation is a one-time implementation cost. OCI tenancy infrastructure consumption is billed "
         "separately at the Core42 / ADGE rate card (procurement action, section 7).",
     size=9.5, color=FAINT, first=True)
footer(s, 6)

# =========================================================
# Slide 7 — Options: cost comparison
# =========================================================
s = slide_new()
header(s, "4", "Options Comparison — Implementation Cost",
       sub="USD, one-time implementation for a comparable 8-module scope. Alternatives are indicative market "
           "benchmark ranges; bar length shows the range midpoint.")
opts = [
    ("A — iFinance on OCI @ Core42", "recommended", 0.095, PLATFORM,
     "USD 100K — fixed quotation  ·  ≈ 3 months"),
    ("B — Extend ADERP Fusion", "custom PaaS extensions via SI", 1.0, NEUTRAL,
     "USD 0.7–1.4M indicative  ·  12–18 months"),
    ("C — COTS point solutions", "multiple products + integration", 0.62, NEUTRAL,
     "USD 0.5–0.9M indicative  ·  6–12 months"),
]
bx, bw = 3.8, 8.9
for i, (label, subl, frac, color, val) in enumerate(opts):
    y = 2.1 + i * 1.15
    tf = tb(s, 0.55, y - 0.08, 3.05, 0.8)
    para(tf, label, size=12.5, bold=True, color=INK, align=PP_ALIGN.RIGHT, first=True, space_after=1)
    para(tf, subl, size=9.5, color=FAINT, align=PP_ALIGN.RIGHT)
    rect(s, bx, y, bw, 0.42, fill=TRACK, round_=True, radius=0.2)
    rect(s, bx, y, max(bw * frac, 0.12), 0.42, fill=color, round_=True, radius=0.2)
    if frac > 0.8:
        tf = tb(s, bx, y, bw - 0.2, 0.42, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, val, size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT, first=True, space_after=0)
    else:
        tf = tb(s, bx + bw * frac + 0.15, y, bw - bw * frac - 0.2, 0.42, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, val, size=11, bold=True, color=INK, first=True, space_after=0)
tf = tb(s, 0.55, 5.9, 12.23, 0.8)
bullets(tf, [
    ("Option A is ~10× lower cost ", "than extending ADERP Fusion, and 4–6× faster than either alternative."),
    ("No recurring application licences ", "— APEX / ORDS / JET runtimes are included with the Oracle database."),
], size=12, gap=6)
footer(s, 7)

# =========================================================
# Slide 8 — Options: criteria table
# =========================================================
s = slide_new()
header(s, "4", "Options Comparison — Evaluation Criteria")
crit = [
    ("Technology alignment",
     "Full Oracle stack, same data centre as ADERP; aligned with ADGE direction",
     "Oracle-aligned, but bound to the shared-instance change pipeline",
     "Mixed vendors; new stacks to secure and support"),
    ("Fusion integration",
     "Both directions, already proven in pilot; in-data-centre traffic",
     "Native, but every change is a central change request",
     "Per-product connectors; typical gaps and duplication"),
    ("Delivery time",
     "≈ 3 months; first go-live week 3, then every two weeks",
     "12–18 months typical",
     "6–12 months typical"),
    ("Cost (indicative)",
     "USD 100K fixed; no per-user licensing",
     "USD 0.7–1.4M",
     "USD 0.5–0.9M + recurring licences"),
    ("Data residency",
     "Abu Dhabi (Core42) — same facility as ADERP",
     "Abu Dhabi (Core42)",
     "Varies by vendor; assessed per product"),
    ("Change speed after go-live",
     "Days — DCT-controlled tenancy and codebase",
     "Weeks–months via the shared change pipeline",
     "Vendor roadmap dependent"),
]
tbl_shape = s.shapes.add_table(len(crit) + 1, 4, Inches(0.55), Inches(1.5), Inches(12.23), Inches(5.1))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(3.7)
tbl.columns[2].width = Inches(3.2)
tbl.columns[3].width = Inches(3.13)
hdr = ["", "A — iFINANCE ON OCI @ CORE42  (RECOMMENDED)", "B — EXTEND ADERP FUSION", "C — COTS POINT SOLUTIONS"]
for j, htxt in enumerate(hdr):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = BRAND if j != 1 else PLATFORM
    c.margin_left = Inches(0.08); c.margin_right = Inches(0.05)
    tf = c.text_frame; tf.word_wrap = True
    para(tf, htxt, size=9.5, bold=True, color=WHITE, first=True, space_after=0)
for i, row in enumerate(crit, start=1):
    fill = WHITE if i % 2 else GROUND
    for j, v in enumerate(row):
        c = tbl.cell(i, j)
        c.fill.solid(); c.fill.fore_color.rgb = SOFT if j == 1 else fill
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.08); c.margin_right = Inches(0.05)
        c.margin_top = c.margin_bottom = Inches(0.03)
        tf = c.text_frame; tf.word_wrap = True
        if j == 0:
            para(tf, v, size=10, bold=True, color=INK, first=True, space_after=0)
        else:
            para(tf, v, size=9.5, color=INK if j == 1 else MUTED, bold=(j == 1), first=True, space_after=0)
footer(s, 8)

# =========================================================
# Slide 9 — Timeline (Gantt)
# =========================================================
s = slide_new()
header(s, "5", "Implementation Timeline — Agile Delivery",
       sub="Tenancy provisioned within one week of contract signature  ·  seven two-week sprints delivering "
           "eight modules  ·  go-live every two weeks from week 3  ·  full scope live by week 15 (≈ 3 months).")
g_rows = [
    ("Tenancy provisioning & platform foundation", 1, 1, PLATFORM, "W1"),
    ("Sprint 1  ·  General Ledger",            2, 2, MODULE,   "W2–3"),
    ("Sprint 2  ·  Reporting & BI",            4, 2, MODULE,   "W4–5"),
    ("Sprint 2  ·  Service Pricing (new build)", 4, 2, MODULE, "W4–5"),
    ("Sprint 3  ·  Freelancers",               6, 2, MODULE,   "W6–7"),
    ("Sprint 4  ·  Team Management",           8, 2, MODULE,   "W8–9"),
    ("Sprint 5  ·  Petty Cash",               10, 2, MODULE,   "W10–11"),
    ("Sprint 6  ·  Duty Travel",              12, 2, MODULE,   "W12–13"),
    ("Sprint 7  ·  Employee Housing (new build)", 14, 2, MODULE, "W14–15"),
    ("Fusion integration · support · hypercare", 2, 14, PLATFORM, "W2–15  (continuous)"),
]
gx0, gy0 = 3.55, 1.95          # chart origin
gw = 9.2                        # chart width
wk = gw / 15.0                  # week width
rh = 0.355                      # row height
# week header
for i in range(15):
    tf = tb(s, gx0 + i * wk, gy0 - 0.28, wk, 0.25)
    para(tf, f"W{i+1}", size=8, color=FAINT, align=PP_ALIGN.CENTER, first=True, space_after=0)
# vertical gridlines
for i in range(16):
    rect(s, gx0 + i * wk, gy0, 0.008, rh * len(g_rows), fill=LINE)
# rows
for r, (label, start, span, color, bl) in enumerate(g_rows):
    y = gy0 + r * rh
    tf = tb(s, 0.55, y + 0.02, 2.95, rh, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, size=9.5, color=INK, first=True, space_after=0, line_spacing=0.95)
    bar_x = gx0 + (start - 1) * wk + 0.02
    bar_w = span * wk - 0.04
    rect(s, bar_x, y + 0.07, bar_w, 0.21, fill=color, round_=True, radius=0.35)
    tf = tb(s, bar_x, y + 0.065, bar_w, 0.21, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, bl, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True, space_after=0)
# milestone rail
rail_y = gy0 + len(g_rows) * rh + 0.12
rect(s, 0.55, rail_y - 0.06, 12.23, 0.015, fill=LINE)
miles = [(1, "M1", "Platform ready", 0), (3, "M2", "First module live (GL)", 1),
         (9, "M3", "5 of 8 modules live", 0), (15, "M4", "Full scope live & handover", 0)]
for wke, mn, ml, lrow in miles:
    mx = gx0 + wke * wk
    d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(mx - 0.07), Inches(rail_y), Inches(0.14), Inches(0.14))
    d.fill.solid(); d.fill.fore_color.rgb = ACCENT; d.line.fill.background(); d.shadow.inherit = False
    if lrow:  # dropped label row: add a faint leader tick from the diamond down to the label
        rect(s, mx - 0.006, rail_y + 0.16, 0.012, 0.30, fill=LINE)
    tf = tb(s, mx - 1.55, rail_y + 0.18 + lrow * 0.34, 1.62, 0.55)
    para(tf, [(mn + "  ", {"bold": True, "color": INK}), (ml, {"color": MUTED})],
         size=8.5, align=PP_ALIGN.RIGHT, first=True, space_after=0, line_spacing=0.95)
# legend
ly = rail_y + 1.0
rect(s, 0.55, ly, 0.18, 0.18, fill=PLATFORM, round_=True, radius=0.25)
tf = tb(s, 0.8, ly - 0.03, 3.4, 0.3)
para(tf, "Platform & integration workstream", size=9.5, color=MUTED, first=True)
rect(s, 4.3, ly, 0.18, 0.18, fill=MODULE, round_=True, radius=0.25)
tf = tb(s, 4.55, ly - 0.03, 4.5, 0.3)
para(tf, "Module delivery sprints (2 weeks each, incl. UAT & training)", size=9.5, color=MUTED, first=True)
d = s.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.45), Inches(ly + 0.02), Inches(0.14), Inches(0.14))
d.fill.solid(); d.fill.fore_color.rgb = ACCENT; d.line.fill.background(); d.shadow.inherit = False
tf = tb(s, 9.68, ly - 0.03, 1.5, 0.3)
para(tf, "Milestone", size=9.5, color=MUTED, first=True)
footer(s, 9)

# =========================================================
# Slide 10 — Inside each sprint
# =========================================================
s = slide_new()
header(s, "5", "Inside Each Two-Week Sprint")
steps = [
    ("DAYS 1–2", "Provision & configure", "Module schema, REST services and organisation data set up on the tenancy."),
    ("DAYS 3–5", "Adapt to DCT policy", "Approval chains, delegation rules, GL mappings and reference data configured."),
    ("DAYS 6–7", "Fusion integration", "Extract feeds and write-back wiring connected and reconciled for the module."),
    ("DAYS 8–9", "UAT & training", "Structured UAT with the business owner; fixes applied; end-user training delivered."),
    ("DAY 10", "Sign-off & go-live", "Sprint review, formal sign-off and production go-live; retrospective feeds the next sprint."),
]
sx, sy, sw_, sh_, sg = 0.55, 2.2, 2.33, 2.9, 0.15
for i, (d, t, body) in enumerate(steps):
    x = sx + i * (sw_ + sg)
    rect(s, x, sy, sw_, sh_, fill=GROUND, line=LINE, round_=True, radius=0.06)
    rect(s, x, sy, sw_, 0.06, fill=ACCENT)
    tf = tb(s, x + 0.2, sy + 0.25, sw_ - 0.4, 0.3)
    para(tf, d, size=10, bold=True, color=ACCENT, first=True)
    tf = tb(s, x + 0.2, sy + 0.58, sw_ - 0.4, 0.65)
    para(tf, t, font=SERIF, size=14, bold=True, color=INK, first=True, line_spacing=1.0)
    tf = tb(s, x + 0.2, sy + 1.25, sw_ - 0.4, sh_ - 1.4)
    para(tf, body, size=10.5, color=MUTED, first=True, line_spacing=1.15)
    if i < 4:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + sw_ + 0.01), Inches(sy + sh_/2 - 0.08),
                                Inches(0.14), Inches(0.16))
        ar.fill.solid(); ar.fill.fore_color.rgb = NEUTRAL; ar.line.fill.background(); ar.shadow.inherit = False
tf = tb(s, 0.55, 5.6, 12.23, 0.5)
para(tf, "Every sprint concludes with a working module in production — accepted by its business owner, "
         "with users trained and a documented UAT evidence pack.",
     size=12, color=MUTED, align=PP_ALIGN.CENTER, first=True)
footer(s, 10)

# =========================================================
# Slide 11 — Governance & quality
# =========================================================
s = slide_new()
header(s, "6", "Governance & Quality")
rect(s, 0.55, 1.5, 0.06, 4.6, fill=BRAND)
tf = tb(s, 0.95, 1.65, 11.7, 4.6)
bullets(tf, [
    ("Agile governance:  ", "sprint planning, sprint review and retrospective conducted with the module's "
     "business owner in every sprint."),
    ("Executive oversight:  ", "fortnightly steering review with the Finance Director, aligned to sprint boundaries."),
    ("Quality gates:  ", "formal sign-off at every sprint close, supported by a documented UAT evidence pack "
     "and delivered training material per module."),
    ("Operational discipline:  ", "versioned deployment runbooks; government security baseline applied at "
     "tenancy provisioning; bilingual English / Arabic experience as a standard acceptance criterion."),
], size=15, gap=16)
footer(s, 11)

# =========================================================
# Slide 12 — Risks & mitigations
# =========================================================
s = slide_new()
header(s, "6", "Risks & Mitigations")
risks = [
    ("Tenancy procurement lead time", "Delays the entire critical path",
     "Raise the Core42 / ADGE procurement immediately upon approval; provisioning completes within one week of contract signature."),
    ("Fusion integration access", "Extract / posting cycles blocked",
     "Request reporting-extract and posting access through ADGE at kick-off; a proven interim extract mechanism from the pilot covers any gap."),
    ("Business availability for UAT & training", "Sprint sign-off slips",
     "Fixed UAT and training window in the second week of each sprint, with a named business owner confirmed at sprint planning."),
    ("New-build scope drift (Service Pricing, Employee Housing)", "A new-build module overruns its sprint",
     "Requirements workshop and scope freeze before Sprint 2 and Sprint 7 respectively; both follow the proven platform pattern."),
    ("Security & residency compliance", "Approval or audit findings",
     "All data resident in Core42 Abu Dhabi; single dedicated tenancy; SSO, role-based access and full audit trail; security review at M1."),
]
tbl_shape = s.shapes.add_table(len(risks) + 1, 3, Inches(0.55), Inches(1.5), Inches(12.23), Inches(4.9))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.3)
tbl.columns[1].width = Inches(2.7)
tbl.columns[2].width = Inches(6.23)
for j, htxt in enumerate(["RISK", "IMPACT", "MITIGATION"]):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = BRAND
    c.margin_left = Inches(0.08)
    tf = c.text_frame; tf.word_wrap = True
    para(tf, htxt, size=10, bold=True, color=WHITE, first=True, space_after=0)
for i, row in enumerate(risks, start=1):
    fill = WHITE if i % 2 else GROUND
    for j, v in enumerate(row):
        c = tbl.cell(i, j)
        c.fill.solid(); c.fill.fore_color.rgb = fill
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_left = Inches(0.08); c.margin_right = Inches(0.05)
        c.margin_top = c.margin_bottom = Inches(0.03)
        tf = c.text_frame; tf.word_wrap = True
        para(tf, v, size=10, bold=(j == 0), color=INK if j == 0 else MUTED, first=True, space_after=0)
footer(s, 12)

# =========================================================
# Slide 13 — Requested approvals & next steps
# =========================================================
s = slide_new()
header(s, "7", "Requested Approvals & Next Steps")
steps13 = [
    ("Approve the recommended option", "and the USD 100,000 implementation budget for the eight-module scope."),
    ("Initiate procurement of the Oracle OCI tenancy", "in the Core42 data centre (same facility as ADERP) — "
     "the single critical-path action; provisioning completes within one week of contract signature."),
    ("Request Fusion integration access", "(reporting extracts and transaction posting) through ADGE, in parallel with procurement."),
    ("Nominate one business owner per module", "to lead UAT and receive training within each sprint."),
    ("Confirm kick-off", "— Sprint 1 (General Ledger) starts in week 2, with the first go-live in week 3."),
]
for i, (lead, rest) in enumerate(steps13):
    y = 1.55 + i * 1.02
    rect(s, 0.55, y, 12.23, 0.88, fill=GROUND, line=LINE, round_=True, radius=0.09)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(y + 0.23), Inches(0.42), Inches(0.42))
    circ.fill.solid(); circ.fill.fore_color.rgb = BRAND; circ.line.fill.background(); circ.shadow.inherit = False
    tfc = circ.text_frame; tfc.word_wrap = False
    para(tfc, str(i + 1), font=SERIF, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, first=True, space_after=0)
    tf = tb(s, 1.55, y + 0.13, 11.0, 0.66, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(lead + " ", {"bold": True, "color": BRAND}), (rest, {"color": MUTED})],
         size=12.5, first=True, space_after=0, line_spacing=1.05)
footer(s, 13)

# ---------------- save ----------------
import os
out_dir = r"c:\claude\DCT-task-management\DCT-Task-Management\docs\presentations"
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "iFinance_Implementation_Plan_08-07-2026.pptx")
prs.save(out)
print("saved:", out)
