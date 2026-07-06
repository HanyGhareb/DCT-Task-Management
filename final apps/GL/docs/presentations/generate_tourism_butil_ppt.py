# -*- coding: utf-8 -*-
"""
Tourism Sector - FY2026 Project Budget Utilization - Executive deck generator.

Data inputs (CSV, produced by ppt_extract*.sql against DCT_BUDGET_UTILIZATION_V
and the reconciling drill fact queries - see extract_tourism_butil.sql):
  c:/tmp/ppt_totals.csv, ppt_by_cc.csv, ppt_mom_ap.csv, ppt_mom_grn.csv,
  ppt_top_projects.csv

Output: Tourism_Budget_Utilization_FY2026_<dd-mm-yyyy>.pptx (+ assets_*/ PNGs)

Chart conventions (dataviz skill): categorical palette validated
  GRN #2F8459 / AP #7C4DBE / Open PR #B0721E / Open PO #3A4FB0
  (Available = neutral track #E3E9E6, not a data series); white 2px gaps
  between stacked segments; recessive grid; text in ink, never series color;
  one axis per chart; legend for >=2 series.
"""
import csv, os, datetime
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import Patch
from matplotlib.ticker import FuncFormatter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------- constants
TMP   = r'c:/tmp'
HERE  = os.path.dirname(os.path.abspath(__file__))
STAMP = '06-07-2026'
ASSET = os.path.join(HERE, 'assets_tourism_fy2026')
OUT   = os.path.join(HERE, 'Tourism_Budget_Utilization_FY2026_%s.pptx' % STAMP)
os.makedirs(ASSET, exist_ok=True)

INK    = '#1F2A27'; MUTED = '#5A6560'; GRID = '#E8EDEA'; TRACK = '#E3E9E6'
C_GRN  = '#2F8459'; C_AP  = '#7C4DBE'; C_PR  = '#B0721E'; C_PO  = '#3A4FB0'
BRAND  = '#3F6F5F'; BRAND_DARK = '#2C5044'; SOFT = '#F2F6F4'
R_INK  = RGBColor(0x1F, 0x2A, 0x27); R_MUTED = RGBColor(0x5A, 0x65, 0x60)
R_BRAND = RGBColor(0x3F, 0x6F, 0x5F); R_BRAND_DARK = RGBColor(0x2C, 0x50, 0x44)
R_WHITE = RGBColor(0xFF, 0xFF, 0xFF); R_SOFT = RGBColor(0xF2, 0xF6, 0xF4)
R_LINE  = RGBColor(0xD8, 0xE0, 0xDC)

plt.rcParams.update({
    'font.family': 'Segoe UI', 'text.color': INK, 'axes.edgecolor': GRID,
    'axes.labelcolor': MUTED, 'xtick.color': MUTED, 'ytick.color': MUTED,
    'axes.linewidth': 0.8, 'figure.facecolor': 'white', 'axes.facecolor': 'white',
    'svg.fonttype': 'none'})

def fm(x, dec=1):
    x = float(x or 0); a = abs(x)
    if a >= 1e9: return '%.2fB' % (x / 1e9)
    if a >= 1e6: return ('%.' + str(dec) + 'fM') % (x / 1e6)
    if a >= 1e3: return '%.0fK' % (x / 1e3)
    return '%.0f' % x

def read_csv(name):
    rows = []
    with open(os.path.join(TMP, name), encoding='utf-8-sig') as f:
        for r in csv.DictReader(f):
            rows.append(r)
    return rows

# ---------------------------------------------------------------- data
tot = read_csv('ppt_totals.csv')[0]
BUDGET = float(tot['BUDGET']); AP = float(tot['ACTUAL_AP']); GRN = float(tot['ACTUAL_GRN'])
PR = float(tot['OPEN_PR']); PO = float(tot['OPEN_PO']); FUND = float(tot['FUND_AVAILABLE'])
ACTUAL = AP + GRN; COMMITTED = PR + PO
UTIL = ACTUAL / BUDGET * 100.0; COMM_PCT = COMMITTED / BUDGET * 100.0

ccs = read_csv('ppt_by_cc.csv')          # sorted by budget desc
for c in ccs:
    for k in ('BUDGET', 'ACTUAL_AP', 'ACTUAL_GRN', 'OPEN_PR', 'OPEN_PO', 'FUND_AVAILABLE'):
        c[k] = float(c[k])
    c['ACTUAL'] = c['ACTUAL_AP'] + c['ACTUAL_GRN']
    c['UTIL'] = c['ACTUAL'] / c['BUDGET'] * 100.0 if c['BUDGET'] else 0.0
    c['COMM'] = c['OPEN_PR'] + c['OPEN_PO']
    c['SHORT'] = (c['DEPARTMENT'].replace('DCT - ', '').replace('DCT-', '')
                  .replace(' Department', '').replace(' Dept.', '').replace(' Dept', ''))
    if c['COST_CENTRE'] == '4515010':
        c['SHORT'] += ' (Office)'          # disambiguate from 4515500 (same dept name)

MONTHS = ['2026-0%d' % m for m in range(1, 8)]
M_LBL  = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul*']
mom_ap, mom_grn = {}, {}
for r in read_csv('ppt_mom_ap.csv'):
    mom_ap[(r['MTH'], r['CC'])] = float(r['AP_AED'])
for r in read_csv('ppt_mom_grn.csv'):
    mom_grn[(r['MTH'], r['CC'])] = float(r['GRN_AED'])
def series(cc=None):
    ap, gr = [], []
    for m in MONTHS:
        if cc: ap.append(mom_ap.get((m, cc), 0.0)); gr.append(mom_grn.get((m, cc), 0.0))
        else:
            ap.append(sum(v for (mm, _), v in mom_ap.items() if mm == m))
            gr.append(sum(v for (mm, _), v in mom_grn.items() if mm == m))
    return ap, gr

top = {}
for r in read_csv('ppt_top_projects.csv'):
    top.setdefault(r['COST_CENTRE'], []).append(r)

# ---------------------------------------------------------------- charts
MFT = FuncFormatter(lambda v, p: fm(v, 0))

def comp_bar(path, b, grn, ap, pr, po, width=8.6, height=1.35, labels=True):
    """Single horizontal composition bar on the full-budget track."""
    fig, ax = plt.subplots(figsize=(width, height), dpi=200)
    segs = [(grn, C_GRN, 'GRN received'), (ap, C_AP, 'AP direct'),
            (pr, C_PR, 'Open PR'), (po, C_PO, 'Open PO')]
    ax.barh([0], [b], color=TRACK, height=.52, zorder=1)
    left = 0.0
    for v, col, _ in segs:
        ax.barh([0], [v], left=left, color=col, height=.52, zorder=2,
                edgecolor='white', linewidth=1.2)
        if labels and v / b > 0.055:
            ax.text(left + v / 2, 0, fm(v, 0), ha='center', va='center',
                    fontsize=8.5, color='white', fontweight='bold', zorder=3)
        left += v
    rem = b - left
    if labels and rem / b > 0.08:
        ax.text(left + rem / 2, 0, fm(rem, 0) + ' available', ha='center',
                va='center', fontsize=8.5, color=MUTED, zorder=3)
    ax.set_xlim(0, b * 1.001); ax.set_ylim(-.55, .55)
    ax.axis('off')
    hnd = [Patch(fc=c, label=l) for _, c, l in segs] + [Patch(fc=TRACK, label='Available')]
    ax.legend(handles=hnd, loc='upper center', bbox_to_anchor=(.5, -.12), ncol=5,
              frameon=False, fontsize=8, handlelength=1.1, handleheight=.9,
              columnspacing=1.2)
    fig.savefig(path, bbox_inches='tight', facecolor='white'); plt.close(fig)

def cc_stack_chart(path):
    """Per cost centre: budget track + Actual/OpenPR/OpenPO stacked, AED."""
    d = ccs; n = len(d)
    fig, ax = plt.subplots(figsize=(8.9, 4.4), dpi=200)
    y = list(range(n))[::-1]
    ax.barh(y, [c['BUDGET'] for c in d], color=TRACK, height=.62, zorder=1)
    left = [0.0] * n
    for key, col, lbl in (('ACTUAL', C_GRN, 'Actual (AP + GRN)'),
                          ('OPEN_PR', C_PR, 'Open PR'), ('OPEN_PO', C_PO, 'Open PO')):
        vals = [c[key] for c in d]
        ax.barh(y, vals, left=left, color=col, height=.62, zorder=2,
                edgecolor='white', linewidth=1.2, label=lbl)
        left = [a + b for a, b in zip(left, vals)]
    for yi, c in zip(y, d):
        ax.text(c['BUDGET'] + BUDGET * .004, yi, fm(c['BUDGET'], 0), va='center',
                fontsize=8, color=MUTED)
    ax.set_yticks(y); ax.set_yticklabels([c['SHORT'] for c in d], fontsize=8.5)
    ax.xaxis.set_major_formatter(MFT); ax.tick_params(axis='x', labelsize=8)
    ax.grid(axis='x', color=GRID, zorder=0); ax.set_axisbelow(True)
    for s in ('top', 'right', 'left'): ax.spines[s].set_visible(False)
    ax.legend(loc='lower right', frameon=False, fontsize=8.5, handlelength=1.1)
    ax.margins(x=.06)
    fig.savefig(path, bbox_inches='tight', facecolor='white'); plt.close(fig)

def cc_pct_chart(path):
    """Per cost centre: % of budget utilised / reserved (PR) / on order (PO)."""
    d = sorted(ccs, key=lambda c: -c['UTIL']); n = len(d)
    y = list(range(n))[::-1]
    fig, ax = plt.subplots(figsize=(8.9, 4.4), dpi=200)
    ax.barh(y, [100] * n, color=TRACK, height=.62, zorder=1)
    left = [0.0] * n
    for key, col, lbl in (('ACTUAL', C_GRN, 'Utilised'), ('OPEN_PR', C_PR, 'Open PR'),
                          ('OPEN_PO', C_PO, 'Open PO')):
        vals = [c[key] / c['BUDGET'] * 100 for c in d]
        ax.barh(y, vals, left=left, color=col, height=.62, zorder=2,
                edgecolor='white', linewidth=1.2, label=lbl)
        left = [a + b for a, b in zip(left, vals)]
    for yi, c in zip(y, d):
        ax.text(1.5, yi, '%.0f%%' % c['UTIL'], va='center', fontsize=8,
                color='white', fontweight='bold', zorder=3)
    ax.axvline(UTIL, color=INK, lw=1.1, ls=(0, (4, 3)), zorder=3)
    ax.text(UTIL + 1, n - .25, 'sector %.0f%%' % UTIL, fontsize=8, color=INK)
    ax.set_yticks(y); ax.set_yticklabels([c['SHORT'] for c in d], fontsize=8.5)
    ax.set_xlim(0, 100); ax.set_xticks([0, 25, 50, 75, 100])
    ax.set_xticklabels(['0%', '25%', '50%', '75%', '100%'], fontsize=8)
    ax.grid(axis='x', color=GRID, zorder=0); ax.set_axisbelow(True)
    for s in ('top', 'right', 'left'): ax.spines[s].set_visible(False)
    ax.legend(loc='lower right', frameon=False, fontsize=8.5, handlelength=1.1)
    fig.savefig(path, bbox_inches='tight', facecolor='white'); plt.close(fig)

def mom_chart(path, cc=None, width=8.9, height=3.9, cum=True):
    ap, gr = series(cc)
    x = list(range(len(MONTHS)))
    fig, ax = plt.subplots(figsize=(width, height), dpi=200)
    ax.bar(x, gr, color=C_GRN, width=.58, label='GRN received',
           edgecolor='white', linewidth=1.2, zorder=2)
    ax.bar(x, ap, bottom=gr, color=C_AP, width=.58, label='AP direct',
           edgecolor='white', linewidth=1.2, zorder=2)
    if cum:
        cumv, s = [], 0.0
        for a, g in zip(ap, gr): s += a + g; cumv.append(s)
        ax.plot(x, cumv, color=INK, lw=2, marker='o', ms=4.5,
                label='Cumulative actual', zorder=3)
        ax.annotate(fm(cumv[-1], 0), (x[-1], cumv[-1]), textcoords='offset points',
                    xytext=(0, 9), ha='center', fontsize=8.5, fontweight='bold', color=INK)
    ax.set_xticks(x); ax.set_xticklabels(M_LBL, fontsize=8.5)
    ax.yaxis.set_major_formatter(MFT); ax.tick_params(axis='y', labelsize=8)
    ax.grid(axis='y', color=GRID, zorder=0); ax.set_axisbelow(True)
    for s in ('top', 'right'): ax.spines[s].set_visible(False)
    ax.legend(loc='upper left', frameon=False, fontsize=8.5, handlelength=1.1)
    ax.margins(y=.12)
    fig.savefig(path, bbox_inches='tight', facecolor='white'); plt.close(fig)

comp_bar(os.path.join(ASSET, 'comp_sector.png'), BUDGET, GRN, AP, PR, PO)
cc_stack_chart(os.path.join(ASSET, 'cc_stack.png'))
cc_pct_chart(os.path.join(ASSET, 'cc_pct.png'))
mom_chart(os.path.join(ASSET, 'mom_sector.png'))
for c in ccs:
    cc = c['COST_CENTRE']
    comp_bar(os.path.join(ASSET, 'comp_%s.png' % cc), c['BUDGET'], c['ACTUAL_GRN'],
             c['ACTUAL_AP'], c['OPEN_PR'], c['OPEN_PO'], width=6.0, height=1.25)
    mom_chart(os.path.join(ASSET, 'mom_%s.png' % cc), cc=cc, width=6.0, height=1.85, cum=False)

# ---------------------------------------------------------------- pptx helpers
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def box(s, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    return sh

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs = list of paragraphs; each = list of (txt, size, bold, color, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (t, size, bold, color, font) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb

SERIF = 'Georgia'; SANS = 'Segoe UI'

def header(s, title, sub=None):
    text(s, 0.55, 0.32, 11.0, 0.65, [[(title, 26, True, R_BRAND_DARK, SERIF)]])
    if sub:
        text(s, 0.55, 0.98, 12.2, 0.35, [[(sub, 12, False, R_MUTED, SANS)]])
    box(s, 0.55, 1.42, 12.23, 0.022, fill=R_BRAND)
    text(s, 11.6, 0.42, 1.7, 0.3, [[('FY2026 · Tourism', 10, True, R_BRAND, SANS)]],
         align=PP_ALIGN.RIGHT)

def kpi_row(s, y, items, x0=0.55, w_total=12.23, h=1.05):
    n = len(items); gap = 0.14; w = (w_total - gap * (n - 1)) / n
    for i, (label, value, note, col) in enumerate(items):
        x = x0 + i * (w + gap)
        box(s, x, y, w, h, fill=R_SOFT)
        box(s, x, y, 0.045, h, fill=col)
        text(s, x + 0.16, y + 0.10, w - 0.25, 0.25, [[(label.upper(), 9, True, R_MUTED, SANS)]])
        text(s, x + 0.16, y + 0.33, w - 0.25, 0.42, [[(value, 19, True, R_INK, SANS)]])
        text(s, x + 0.16, y + 0.76, w - 0.25, 0.24, [[(note, 8.5, False, R_MUTED, SANS)]])

def pic(s, path, x, y, w):
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def bullets(s, x, y, w, h, items, size=12.5, gap=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = '▪  '; r.font.size = Pt(size); r.font.color.rgb = R_BRAND
        r.font.name = SANS
        r = p.add_run(); r.text = head + ' — '; r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = R_INK; r.font.name = SANS
        r = p.add_run(); r.text = body; r.font.size = Pt(size); r.font.bold = False
        r.font.color.rgb = R_INK; r.font.name = SANS

# ---------------------------------------------------------------- S1 title
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=R_BRAND_DARK)
box(s, 0, 4.62, 13.333, 0.03, fill=RGBColor(0xC8, 0x9B, 0x3C))
text(s, 1.0, 1.55, 11.3, 0.5, [[('DEPARTMENT OF CULTURE AND TOURISM · FINANCE — i-FINANCE GL', 13, True,
                                 RGBColor(0xBF, 0xD5, 0xCB), SANS)]])
text(s, 1.0, 2.15, 11.3, 1.9, [
    [('Tourism Sector', 48, True, R_WHITE, SERIF)],
    [('Project Budget Utilization — FY2026', 30, False, RGBColor(0xDD, 0xE9, 0xE2), SERIF)]])
text(s, 1.0, 4.85, 11.3, 1.4, [
    [('Executive briefing — budget vs. execution, open commitments and month-over-month delivery', 14, False, RGBColor(0xC9, 0xDA, 0xD1), SANS)],
    [('Prepared 6 July 2026 · Source: i-Finance General Ledger — Budget Utilization (live)', 11.5, False, RGBColor(0x9F, 0xB8, 0xAC), SANS)]])

# ---------------------------------------------------------------- S2 executive summary
s = slide()
header(s, 'Executive summary', 'FY2026 project budget position — all figures AED, year-to-date through 6 July 2026')
kpi_row(s, 1.75, [
    ('Approved budget', fm(BUDGET), '140 projects · 8 cost centres', R_BRAND),
    ('Actual spend', fm(ACTUAL), '%.1f%% utilised (GRN %s + AP %s)' % (UTIL, fm(GRN, 0), fm(AP, 0)), RGBColor(0x2F, 0x84, 0x59)),
    ('Open commitments (PR)', fm(PR), 'reserved requisitions', RGBColor(0xB0, 0x72, 0x1E)),
    ('Open obligations (PO)', fm(PO), 'orders awaiting delivery', RGBColor(0x3A, 0x4F, 0xB0)),
    ('Funds available', fm(FUND), '%.1f%% of budget uncommitted' % (FUND / BUDGET * 100), R_MUTED)])
text(s, 0.55, 3.12, 12.0, 0.3, [[('How the budget sits today', 14, True, R_INK, SANS)]])
pic(s, os.path.join(ASSET, 'comp_sector.png'), 0.95, 3.45, 11.4)
bullets(s, 0.55, 6.02, 12.2, 1.4, [
    ('60% of the budget is already deployed or spoken for',
     'actual %s (%.1f%%) plus %s of open PR/PO commitments leave %s genuinely free.' %
     (fm(ACTUAL), UTIL, fm(COMMITTED), fm(FUND))),
    ('Execution is procurement-led', 'over 99% of actual spend arrives through goods/services receipts (GRN) rather than direct invoices — delivery follow-through drives utilisation.'),
    ('Concentration', 'Tourism (Central) holds %.0f%% of the sector budget; its execution pattern defines the sector curve.' % (ccs[0]['BUDGET'] / BUDGET * 100))],
    size=11.5, gap=4)

# ---------------------------------------------------------------- S3 by CC (AED)
s = slide()
header(s, 'Budget by cost centre', 'Approved budget (track) vs. actual spend and open commitments — AED')
pic(s, os.path.join(ASSET, 'cc_stack.png'), 0.55, 1.75, 12.2)
text(s, 0.55, 6.95, 12.2, 0.35, [[('Bars show each cost centre’s full budget (grey track); coloured segments are actual spend and open PR/PO within it. Budget value at bar end.',
                                   9.5, False, R_MUTED, SANS)]])

# ---------------------------------------------------------------- S4 by CC (%)
s = slide()
header(s, 'Utilisation ranking', 'Share of each cost centre’s own budget — utilised, reserved (PR) and on order (PO)')
pic(s, os.path.join(ASSET, 'cc_pct.png'), 0.55, 1.75, 12.2)
text(s, 0.55, 6.95, 12.2, 0.35, [[('Dashed line = sector utilisation (%.1f%%). A low utilised share with a large amber/indigo block means the pipeline is committed but not yet delivered.' % UTIL,
                                   9.5, False, R_MUTED, SANS)]])

# ---------------------------------------------------------------- S5 MoM
s = slide()
header(s, 'Month-over-month execution', 'Monthly actual spend (GRN + AP direct) with cumulative curve — AED')
pic(s, os.path.join(ASSET, 'mom_sector.png'), 0.9, 1.8, 11.5)
text(s, 0.55, 6.95, 12.2, 0.35, [[('* July is month-to-date (through 6 Jul). March reflects a major receipt cycle at Tourism (Central) — execution lands in steps as deliveries are receipted, not in a straight line.',
                                   9.5, False, R_MUTED, SANS)]])

# ---------------------------------------------------------------- S6.. deep dives
for c in ccs:
    cc = c['COST_CENTRE']
    s = slide()
    header(s, c['SHORT'], 'Cost centre %s · %d projects · deep dive' % (cc, int(c['PROJECTS_N'])))
    kpi_row(s, 1.62, [
        ('Budget', fm(c['BUDGET']), '%.1f%% of sector' % (c['BUDGET'] / BUDGET * 100), R_BRAND),
        ('Actual', fm(c['ACTUAL']), '%.1f%% utilised' % c['UTIL'], RGBColor(0x2F, 0x84, 0x59)),
        ('Open PR', fm(c['OPEN_PR']), '%.1f%% of budget' % (c['OPEN_PR'] / c['BUDGET'] * 100), RGBColor(0xB0, 0x72, 0x1E)),
        ('Open PO', fm(c['OPEN_PO']), '%.1f%% of budget' % (c['OPEN_PO'] / c['BUDGET'] * 100), RGBColor(0x3A, 0x4F, 0xB0)),
        ('Available', fm(c['FUND_AVAILABLE']), '%.1f%% of budget' % (c['FUND_AVAILABLE'] / c['BUDGET'] * 100), R_MUTED)], h=0.98)
    text(s, 0.55, 2.78, 6.1, 0.28, [[('Budget composition', 12, True, R_INK, SANS)]])
    pic(s, os.path.join(ASSET, 'comp_%s.png' % cc), 0.55, 3.08, 6.1)
    text(s, 7.0, 2.78, 5.8, 0.28, [[('Monthly execution (AED)', 12, True, R_INK, SANS)]])
    pic(s, os.path.join(ASSET, 'mom_%s.png' % cc), 7.0, 3.08, 5.8)
    # top projects table
    text(s, 0.55, 5.02, 8.0, 0.28, [[('Top projects by budget', 12, True, R_INK, SANS)]])
    rows = top.get(cc, [])[:5]
    tbl = s.shapes.add_table(len(rows) + 1, 5, Inches(0.55), Inches(5.32),
                             Inches(12.2), Inches(0.32 * (len(rows) + 1))).table
    tbl.columns[0].width = Inches(4.9); tbl.columns[1].width = Inches(1.9)
    tbl.columns[2].width = Inches(1.9); tbl.columns[3].width = Inches(1.9)
    tbl.columns[4].width = Inches(1.6)
    hdrs = ['Project', 'Budget', 'Actual', 'Committed (PR+PO)', 'Available']
    for j, htxt in enumerate(hdrs):
        cell = tbl.cell(0, j); cell.fill.solid(); cell.fill.fore_color.rgb = R_BRAND
        cell.margin_top = cell.margin_bottom = Pt(2)
        p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = R_WHITE; r.font.name = SANS
        if j: p.alignment = PP_ALIGN.RIGHT
    for i, pr_row in enumerate(rows, start=1):
        vals = ['%s — %s' % (pr_row['PROJECT_NUMBER'], pr_row['PROJECT_NAME']),
                fm(float(pr_row['BUDGET'])), fm(float(pr_row['ACTUAL'])),
                fm(float(pr_row['COMMITTED'])), fm(float(pr_row['FUND_AVAILABLE']))]
        for j, v in enumerate(vals):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = R_SOFT if i % 2 else R_WHITE
            cell.margin_top = cell.margin_bottom = Pt(1)
            p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = v
            r.font.size = Pt(9.5); r.font.color.rgb = R_INK; r.font.name = SANS
            if j: p.alignment = PP_ALIGN.RIGHT

# ---------------------------------------------------------------- S14 observations
s = slide()
header(s, 'Observations & recommended actions', 'What the numbers say — and where management attention pays off')
cc5500 = next(c for c in ccs if c['COST_CENTRE'] == '4515500')
cc5300 = next(c for c in ccs if c['COST_CENTRE'] == '4515300')
run_rate = ACTUAL / 6.2
bullets(s, 0.55, 1.75, 12.2, 5.4, [
    ('Concentration risk', 'Tourism (Central) carries %.0f%% of the sector budget (%s); sector delivery is effectively this cost centre’s delivery. Recommend a dedicated monthly delivery review for its top-5 projects.' %
     (cc5500['BUDGET'] / BUDGET * 100, fm(cc5500['BUDGET']))),
    ('Committed but not delivered = %s (%.0f%% of budget)' % (fm(COMMITTED), COMM_PCT),
     'PO %s awaiting receipts plus PR %s awaiting conversion. Utilisation will step up mechanically as deliveries land — push GRN discipline (receipt on delivery, not at invoice time).' % (fm(PO), fm(PR))),
    ('Pipeline stuck at requisition', 'Visitor Experience & Tourism Products Development is only %.0f%% utilised, yet %.0f%% of its budget (%s) is already reserved on requisitions — expedite PR→PO conversion and delivery scheduling.' %
     (cc5300['UTIL'], cc5300['OPEN_PR'] / cc5300['BUDGET'] * 100, fm(cc5300['OPEN_PR']))),
    ('Run-rate check', 'actual spend averages ≈ %s/month. With %s free and %s committed, full-year consumption projects close to budget only if committed orders deliver in H2 — track PO ageing monthly.' %
     (fm(run_rate), fm(FUND), fm(COMMITTED))),
    ('Low-utilisation watchlist', 'Tourism Central Office (8.0%) and Emiratization (7.5%) — small budgets, but if unspent by Q3 they are reallocation candidates.'),
    ('Healthy signal', 'only %s (0.3%%) of spend bypasses procurement as direct AP — strong contract/PO discipline across the sector.' % fm(AP))],
    size=12.5, gap=8)

# ---------------------------------------------------------------- S15 appendix
s = slide()
header(s, 'Appendix — definitions & method', 'How each figure is derived (all AED, FY2026)')
bullets(s, 0.55, 1.8, 12.2, 4.6, [
    ('Approved budget', 'FY2026 project budget per Project × Task × Expenditure Type (Fusion projects budget extract).'),
    ('Actual — GRN received', 'value of goods/services receipted against purchase orders (receipt date in FY2026).'),
    ('Actual — AP direct', 'validated supplier invoices charged directly to a project with no purchase order.'),
    ('Open commitments (PR)', 'requisition lines with funds status Reserved — budget held, order not yet placed.'),
    ('Open obligations (PO)', 'purchase-order lines Reserved / Partially Liquidated, net of receipts to date.'),
    ('Funds available', 'Budget − Actual AP − Actual GRN − Open PR − Open PO.'),
    ('Source', 'i-Finance General Ledger (App 210) — Budget Utilization report over DCT_BUDGET_UTILIZATION_V; monthly split reconciles to the report’s drill-down. Data as of 6 July 2026; July is month-to-date.')],
    size=12, gap=7)

prs.save(OUT)
print('SAVED:', OUT)
