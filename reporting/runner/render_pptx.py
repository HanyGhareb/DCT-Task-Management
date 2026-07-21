"""render_pptx.py — executive PowerPoint decks for the GL reporting books.

build_deck(sections, ctx) -> .pptx bytes. Consumes a MULTI report's rendered
sections and produces a leadership-facing 16:9 deck with NATIVE PowerPoint
tables and charts (fully editable in PowerPoint), so every figure matches the
PDF/Excel outputs. Dispatches on ctx["report_code"]:

  BUDGET_UTIL_BOOK  → the Budget Utilization deck (overview / by_sector /
                      pressure / ap_lines / grn_lines / open_po / open_pr):
                      cover · executive KPIs · utilization by sector · budget
                      composition · pressure lines · actuals & top suppliers ·
                      open obligations & commitments · insights · methodology.
  ENC_PENDING_BOOK  → the Encumbrances Pending Approval deck (overview /
                      budget_ctx / aging / approvers / sectors / oldest /
                      pr_register / po_register / unmatched): cover ·
                      pending-approval overview · aging analysis · pending
                      value by sector · approver follow-up · longest-waiting
                      documents · insights · methodology.

Design language mirrors the Briefing Books: deep teal + gold on white, prepared
by Financial Planning and Budgeting, Finance Department.
"""
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

# ── palette ──────────────────────────────────────────────────────────────────
TEAL = RGBColor(0x1F, 0x6F, 0x8B)
TEAL_D = RGBColor(0x14, 0x4E, 0x63)
GOLD = RGBColor(0xC8, 0xA2, 0x4B)
INK = RGBColor(0x21, 0x2B, 0x30)
MUTE = RGBColor(0x64, 0x74, 0x7A)
LINE = RGBColor(0xDD, 0xE5, 0xE8)
SOFT = RGBColor(0xEE, 0xF3, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OK = RGBColor(0x1E, 0x5E, 0x2C)
BAD = RGBColor(0xB3, 0x26, 0x1E)
# categorical ramp for charts (composition / sectors)
RAMP = [RGBColor(0x1F, 0x6F, 0x8B), RGBColor(0xC8, 0xA2, 0x4B), RGBColor(0x4E, 0x9C, 0xB0),
        RGBColor(0x7A, 0x5C, 0x2E), RGBColor(0x9C, 0xC3, 0xCE), RGBColor(0x2E, 0x7D, 0x5B),
        RGBColor(0xB9, 0x74, 0x1C), RGBColor(0x86, 0xA8, 0xB2)]

EMU_IN = 914400
SW = Inches(13.333)
SH = Inches(7.5)


# ── helpers ──────────────────────────────────────────────────────────────────
def _num(v):
    try:
        if v is None:
            return 0.0
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def money(v):
    """B/M/K scaled, matching the Briefing Book money() macro."""
    v = _num(v)
    a = abs(v)
    if a >= 1e9:
        return f"{v/1e9:,.2f} B"
    if a >= 1e6:
        return f"{v/1e6:,.1f} M"
    if a >= 1e3:
        return f"{v/1e3:,.0f} K"
    return f"{v:,.0f}"


def num(v, dp=0):
    return f"{_num(v):,.{dp}f}"


def _d(v):
    """Format a date/datetime (or already-string date) as YYYY-MM-DD."""
    if v is None:
        return ""
    try:
        return v.strftime("%Y-%m-%d")
    except AttributeError:
        return str(v)[:10]


def _sections_by_key(sections):
    out = {}
    for s in (sections or []):
        out[s.get("key")] = s
    return out


def _data(sec):
    return (sec or {}).get("data") or []


def _txt(shapes, left, top, width, height, text, size=18, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False, font="Public Sans"):
    box = shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    r.font.color.rgb = color
    return box


def _rect(shapes, left, top, width, height, fill, line=None):
    sh = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _band(slide, part_no, title, brand=TEAL):
    """Top brand band with a part number + title (content slides)."""
    _rect(slide.shapes, 0, 0, SW, Inches(0.9), brand)
    _rect(slide.shapes, 0, Inches(0.9), SW, Emu(int(0.05 * EMU_IN)), GOLD)
    if part_no:
        _txt(slide.shapes, Inches(0.5), Inches(0.12), Inches(1.0), Inches(0.66),
             part_no, size=26, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
    _txt(slide.shapes, Inches(1.5), Inches(0.12), Inches(11.3), Inches(0.66),
         title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


def _foot(slide, meta):
    if not meta:                        # footer suppressed (see build_deck)
        return
    _txt(slide.shapes, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
         meta, size=9, color=MUTE, anchor=MSO_ANCHOR.MIDDLE)


def _style_table(table, header_fill=TEAL, header_fg=WHITE, font_size=11, header_size=11):
    """Flatten python-pptx's default banded style + apply the brand look."""
    tbl = table._tbl
    # strip the default table style (removes the blue banding)
    for el in tbl.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tableStyleId'):
        el.getparent().remove(el)
    ncols = len(table.columns)
    for ri, row in enumerate(table.rows):
        for ci in range(ncols):
            cell = table.cell(ri, ci)
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(2)
            cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = cell.text_frame.paragraphs[0]
            if not para.runs:
                para.add_run()
            run = para.runs[0]
            run.font.size = Pt(header_size if ri == 0 else font_size)
            run.font.name = "Public Sans"
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.color.rgb = header_fg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else SOFT
                run.font.color.rgb = INK


def _add_table(slide, left, top, width, headers, rows, col_w=None,
               aligns=None, height=None):
    nrows = len(rows) + 1
    ncols = len(headers)
    height = height or Inches(0.3 * nrows)
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    table = shape.table
    if col_w:
        total = sum(col_w)
        for ci, w in enumerate(col_w):
            table.columns[ci].width = Emu(int(width * w / total))
    for ci, h in enumerate(headers):
        table.cell(0, ci).text = str(h)
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            table.cell(ri, ci).text = "" if val is None else str(val)
    _style_table(table)
    if aligns:
        for ri in range(nrows):
            for ci, al in enumerate(aligns):
                table.cell(ri, ci).text_frame.paragraphs[0].alignment = al
    return table


def _chart(slide, chart_type, left, top, width, height, categories, series,
           colors=None, legend=True, data_labels=False, number_format=None):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals, number_format=number_format or '#,##0')
    gframe = slide.shapes.add_chart(chart_type, left, top, width, height, cd)
    chart = gframe.chart
    chart.has_title = False
    try:
        chart.font.size = Pt(10)
        chart.font.name = "Public Sans"
    except Exception:
        pass
    if legend and len(series) > 1 or (chart_type in (XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.PIE)):
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False
    # colour the points (single-series bar) or the slices (pie/doughnut)
    if colors:
        plot = chart.plots[0]
        for si, ser in enumerate(chart.series):
            if chart_type in (XL_CHART_TYPE.DOUGHNUT, XL_CHART_TYPE.PIE):
                for pi, point in enumerate(ser.points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = colors[pi % len(colors)]
            else:
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = colors[si % len(colors)]
    if data_labels:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = number_format or '#,##0'
        plot.data_labels.number_format_is_linked = False
        plot.data_labels.font.size = Pt(9)
    return chart


def _scope_line(params):
    bits = []
    p = params or {}
    if p.get("period"):
        bits.append(f"YTD through {p['period']}")
    if p.get("sector"):
        bits.append(f"Sector: {p['sector']}")
    if p.get("chapter"):
        bits.append(f"Chapter: {p['chapter']}")
    if p.get("projecttype"):
        bits.append(f"Project type: {p['projecttype']}")
    if p.get("costcenter"):
        bits.append(f"Cost centre: {p['costcenter']}")
    if p.get("project"):
        bits.append(f"Project: {p['project']}")
    if p.get("task"):
        bits.append(f"Task: {p['task']}")
    if p.get("etype"):
        bits.append(f"Expenditure type: {p['etype']}")
    if p.get("search"):
        bits.append(f"Search: {p['search']}")
    return " · ".join(bits) if bits else "Whole organisation · Full year"


# ── Scope slide (Part 1) — the full report filter set the deck was run with ────
_SCOPE_LABELS = [
    ("year", "Budget year"),
    ("period", "Period (YTD)"),
    ("source", "Source"),
    ("sector", "Sector"),
    ("chapter", "Chapter"),
    ("projecttype", "Project type"),
    ("bu", "Business unit(s)"),
    ("costcenter", "Cost centre"),
    ("project", "Project"),
    ("task", "Task"),
    ("etype", "Expenditure type"),
    ("search", "Search text"),
]


def _scope_value(key, raw):
    """Human display of a selected parameter, or its 'all'/default when unset."""
    if raw is None or str(raw) == "":
        return {"year": "—", "period": "Full year",
                "source": "PR & PO", "search": "—"}.get(key, "All")
    v = str(raw)
    if "|" in v:                                   # pipe-delimited any-of list
        v = ", ".join(part for part in v.split("|") if part)
    if key == "period":
        v = "YTD through " + v
    return v if len(v) <= 46 else v[:44] + "…"


def _scope_pairs(ctx):
    p = ctx.get("params") or {}
    code = (ctx.get("report_code") or "").upper()
    pairs = []
    for key, label in _SCOPE_LABELS:
        if key == "source" and code != "ENC_PENDING_BOOK":
            continue                               # 'source' is pending-only
        pairs.append((label, _scope_value(key, p.get(key))))
    return pairs


def _scope_slide(prs, ctx, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "01", "Scope")
    _txt(slide.shapes, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.45),
         "The parameters this report was generated with — every figure in this deck reflects this exact selection.",
         size=12, color=MUTE)
    pairs = _scope_pairs(ctx)
    half = (len(pairs) + 1) // 2
    cols = [pairs[:half], pairs[half:]]
    col_x = [Inches(0.5), Inches(6.78)]
    col_w = Inches(6.05)
    row_h = 0.86
    for ci, col in enumerate(cols):
        y = 1.75
        for label, value in col:
            _rect(slide.shapes, col_x[ci], Inches(y), col_w, Inches(row_h - 0.14), WHITE, line=LINE)
            _rect(slide.shapes, col_x[ci], Inches(y), Emu(int(0.06 * EMU_IN)), Inches(row_h - 0.14), GOLD)
            _txt(slide.shapes, col_x[ci] + Inches(0.22), Inches(y + 0.07), col_w - Inches(0.4), Inches(0.3),
                 label.upper(), size=10, bold=True, color=MUTE)
            _txt(slide.shapes, col_x[ci] + Inches(0.22), Inches(y + 0.33), col_w - Inches(0.4), Inches(0.36),
                 value, size=15, bold=True, color=TEAL_D, anchor=MSO_ANCHOR.MIDDLE)
            y += row_h
    _foot(slide, meta)


# ── slides ───────────────────────────────────────────────────────────────────
def _cover(prs, ctx, O):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide.shapes, 0, 0, SW, Inches(2.35), TEAL)
    _rect(slide.shapes, 0, Inches(2.35), SW, Emu(int(0.08 * EMU_IN)), GOLD)
    _rect(slide.shapes, 0, Inches(7.18), SW, Inches(0.32), TEAL_D)
    p = ctx.get("params") or {}
    _txt(slide.shapes, Inches(0.7), Inches(0.55), Inches(12), Inches(0.4),
         "i-Finance · Reporting Platform", size=14, bold=True, color=RGBColor(0xCF, 0xE3, 0xEA))
    _txt(slide.shapes, Inches(0.7), Inches(1.0), Inches(12), Inches(1.2),
         "Budget Utilization Briefing", size=44, bold=True, color=WHITE)
    sub = f"Budget Year {p.get('year','')}"
    if p.get("period"):
        sub += f"  —  YTD through {p['period']}"
    _txt(slide.shapes, Inches(0.7), Inches(2.6), Inches(12), Inches(0.5),
         sub, size=20, bold=True, color=TEAL_D)
    _txt(slide.shapes, Inches(0.7), Inches(3.25), Inches(12), Inches(0.5),
         "Scope:  " + _scope_line(p), size=13, color=MUTE)
    # KPI ribbon on the cover
    if O:
        tiles = [
            ("Approved Budget", money(O.get("budget"))),
            ("Total Actual", money(_num(O.get("actual_ap")) + _num(O.get("actual_grn")))),
            ("Total Encumbrance", money(_num(O.get("obligation_po")) + _num(O.get("commitment_pr")))),
            ("Fund Available", money(O.get("fund_available"))),
            ("Utilization", (num(O.get("utilization_pct"), 1) + "%") if O.get("utilization_pct") is not None else "—"),
        ]
        n = len(tiles)
        gap = Inches(0.2)
        tw = (SW - Inches(1.4) - gap * (n - 1)) / n
        x = Inches(0.7)
        for i, (lbl, val) in enumerate(tiles):
            _rect(slide.shapes, x, Inches(4.15), tw, Inches(1.35), SOFT, line=LINE)
            _rect(slide.shapes, x, Inches(4.15), tw, Emu(int(0.06 * EMU_IN)), GOLD if i == 4 else TEAL)
            _txt(slide.shapes, x, Inches(4.32), tw, Inches(0.35),
                 lbl.upper(), size=10, bold=True, color=MUTE, align=PP_ALIGN.CENTER)
            _txt(slide.shapes, x, Inches(4.7), tw, Inches(0.7),
                 val, size=22, bold=True, color=TEAL_D, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
            x += tw + gap
    _txt(slide.shapes, Inches(0.7), Inches(6.1), Inches(9), Inches(0.7),
         "Prepared by Financial Planning and Budgeting\nFinance Department",
         size=13, bold=True, color=INK)
    _txt(slide.shapes, Inches(0.7), Inches(7.2), Inches(12), Inches(0.28),
         f"Generated {ctx.get('generated_at','')} GST (UTC+04:00), Asia/Dubai   ·   {ctx.get('report_code','')}",
         size=9, color=RGBColor(0xCF, 0xE3, 0xEA), anchor=MSO_ANCHOR.MIDDLE)


def _kpi_slide(prs, ctx, O, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "01", "Executive Overview")
    if not O:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No approved budget in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    budget = _num(O.get("budget"))
    ap = _num(O.get("actual_ap"))
    grn = _num(O.get("actual_grn"))
    po = _num(O.get("obligation_po"))
    pr = _num(O.get("commitment_pr"))
    fund = _num(O.get("fund_available"))
    util = O.get("utilization_pct")
    tiles = [
        ("Approved Budget", money(budget), f"{num(O.get('budget_lines'))} budget lines", TEAL),
        ("Total Actual (AP+GRN)", money(ap + grn), f"AP {money(ap)} · GRN {money(grn)}", RGBColor(0x6C, 0x4A, 0xB6)),
        ("Total Encumbrance (PR+PO)", money(pr + po), f"PR {money(pr)} · PO {money(po)}", GOLD),
        ("Fund Available", money(fund), (num(100 * fund / budget, 1) + "% of budget" if budget else ""),
         OK if fund >= 0 else BAD),
    ]
    n = len(tiles)
    gap = Inches(0.25)
    tw = (SW - Inches(1.0) - gap * (n - 1)) / n
    x = Inches(0.5)
    for lbl, val, sub, accent in tiles:
        _rect(slide.shapes, x, Inches(1.3), tw, Inches(1.9), WHITE, line=LINE)
        _rect(slide.shapes, x, Inches(1.3), Emu(int(0.08 * EMU_IN)), Inches(1.9), accent)
        _txt(slide.shapes, x + Inches(0.2), Inches(1.5), tw - Inches(0.3), Inches(0.55),
             lbl.upper(), size=11, bold=True, color=MUTE)
        _txt(slide.shapes, x + Inches(0.2), Inches(2.05), tw - Inches(0.3), Inches(0.7),
             val, size=30, bold=True, color=INK)
        _txt(slide.shapes, x + Inches(0.2), Inches(2.75), tw - Inches(0.3), Inches(0.4),
             sub, size=11, color=MUTE)
        x += tw + gap
    # utilization bar
    _txt(slide.shapes, Inches(0.5), Inches(3.55), Inches(6), Inches(0.4),
         "Committed utilization of budget", size=13, bold=True, color=INK)
    upct = _num(util)
    bar_w = SW - Inches(1.0)
    _rect(slide.shapes, Inches(0.5), Inches(4.0), bar_w, Inches(0.5), SOFT, line=LINE)
    filled = max(0.0, min(1.0, upct / 100.0))
    if filled > 0:
        _rect(slide.shapes, Inches(0.5), Inches(4.0), Emu(int(bar_w * filled)), Inches(0.5),
              BAD if upct > 100 else TEAL)
    _txt(slide.shapes, Inches(0.5), Inches(4.05), bar_w - Inches(0.2), Inches(0.4),
         (num(util, 1) + "% utilized" if util is not None else "—"),
         size=13, bold=True, color=WHITE if filled > 0.12 else INK, anchor=MSO_ANCHOR.MIDDLE)
    # composition doughnut
    comp = [("Actual AP", ap), ("Actual GRN", grn), ("Open PO", po), ("Open PR", pr),
            ("Fund Available", max(fund, 0))]
    comp = [(k, v) for k, v in comp if v > 0]
    if comp:
        _txt(slide.shapes, Inches(0.5), Inches(4.85), Inches(6), Inches(0.4),
             "Budget composition", size=13, bold=True, color=INK)
        _chart(slide, XL_CHART_TYPE.DOUGHNUT, Inches(0.4), Inches(5.15), Inches(6.2), Inches(1.85),
               [c[0] for c in comp], [("Composition", [round(c[1], 2) for c in comp])],
               colors=RAMP, legend=True)
    # over-budget callout
    over = int(_num(O.get("over_budget_lines")))
    _rect(slide.shapes, Inches(7.0), Inches(5.15), Inches(5.85), Inches(1.85),
          SOFT, line=LINE)
    _txt(slide.shapes, Inches(7.2), Inches(5.3), Inches(5.5), Inches(0.4),
         "Budget pressure", size=13, bold=True, color=INK)
    _txt(slide.shapes, Inches(7.2), Inches(5.7), Inches(5.5), Inches(0.8),
         str(over), size=40, bold=True, color=BAD if over > 0 else OK)
    _txt(slide.shapes, Inches(7.2), Inches(6.55), Inches(5.5), Inches(0.4),
         "budget line(s) over their approved budget" if over else "every line still has funds available",
         size=12, color=MUTE)
    _foot(slide, meta)


def _sector_slide(prs, ctx, SEC, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "02", "Utilization by Sector")
    rows = sorted(SEC, key=lambda r: -_num(r.get("budget")))[:8]
    if not rows:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No budget in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    # bar chart: Budget vs Committed (actual + encumbrance) per sector
    cats = [str(r.get("sector") or "—")[:22] for r in rows]
    budget = [round(_num(r.get("budget")) / 1e6, 2) for r in rows]
    committed = [round((_num(r.get("actual_ap")) + _num(r.get("actual_grn")) +
                        _num(r.get("obligation_po")) + _num(r.get("commitment_pr"))) / 1e6, 2) for r in rows]
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(1.2),
           Inches(12.5), Inches(3.0), cats,
           [("Budget (AED M)", budget), ("Committed (AED M)", committed)],
           colors=[TEAL, GOLD], legend=True)
    # table
    hdr = ["Sector", "Budget", "Actual", "Encumbrance", "Fund Avail.", "Util %"]
    trows = []
    for r in rows:
        act = _num(r.get("actual_ap")) + _num(r.get("actual_grn"))
        enc = _num(r.get("obligation_po")) + _num(r.get("commitment_pr"))
        trows.append([str(r.get("sector") or "—"), money(r.get("budget")), money(act),
                      money(enc), money(r.get("fund_available")),
                      (num(r.get("utilization_pct"), 1) + "%") if r.get("utilization_pct") is not None else "—"])
    aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 5
    _add_table(slide, Inches(0.4), Inches(4.4), Inches(12.5), hdr, trows,
               col_w=[3, 1.4, 1.4, 1.4, 1.4, 1], aligns=aligns, height=Inches(2.4))
    _foot(slide, meta)


def _pressure_slide(prs, ctx, PRS, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "03", "Budget Lines under Pressure")
    rows = PRS[:12]
    if not rows:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No budget lines in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    _txt(slide.shapes, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         "The lines with the least remaining funds (lowest Fund Available first).",
         size=12, color=MUTE)
    hdr = ["Project", "Task", "Expenditure type", "Budget", "Actual", "Encumbrance", "Fund Avail.", "Util %"]
    trows = []
    for r in rows:
        enc = _num(r.get("obligation_po")) + _num(r.get("commitment_pr"))
        trows.append([
            (str(r.get("project_number") or "") + " " + str(r.get("project_name") or ""))[:34],
            str(r.get("task_number") or "")[:16],
            str(r.get("expenditure_type") or "")[:26],
            money(r.get("budget")), money(r.get("actual")), money(enc),
            money(r.get("fund_available")),
            (num(r.get("utilization_pct"), 1) + "%") if r.get("utilization_pct") is not None else "—"])
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 5
    _add_table(slide, Inches(0.4), Inches(1.7), Inches(12.55), hdr, trows,
               col_w=[3, 1.5, 2.4, 1.2, 1.2, 1.3, 1.3, 0.9], aligns=aligns, height=Inches(5.0))
    _foot(slide, meta)


def _supplier_totals(AP, GRN):
    sup = {}
    for r in AP:
        k = r.get("supplier_name") or "(Unknown supplier)"
        sup[k] = sup.get(k, 0.0) + _num(r.get("matched_aed"))
    for r in GRN:
        k = r.get("supplier_name") or "(Unknown supplier)"
        sup[k] = sup.get(k, 0.0) + _num(r.get("received_aed"))
    return sorted(sup.items(), key=lambda kv: -kv[1])


def _actuals_slide(prs, ctx, AP, GRN, O, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "04", "Actuals & Supplier Concentration")
    ap_tot = sum(_num(r.get("matched_aed")) for r in AP)
    grn_tot = sum(_num(r.get("received_aed")) for r in GRN)
    # left: AP vs GRN mini tiles
    for i, (lbl, val, sub, accent) in enumerate([
        ("Actual — AP (direct invoices)", money(ap_tot), f"{len(AP):,} invoices", TEAL),
        ("Actual — GRN (goods received)", money(grn_tot), f"{len(GRN):,} receipt lines", GOLD)]):
        y = Inches(1.35 + i * 1.35)
        _rect(slide.shapes, Inches(0.5), y, Inches(5.6), Inches(1.15), WHITE, line=LINE)
        _rect(slide.shapes, Inches(0.5), y, Emu(int(0.08 * EMU_IN)), Inches(1.15), accent)
        _txt(slide.shapes, Inches(0.75), y + Inches(0.12), Inches(5.2), Inches(0.4),
             lbl.upper(), size=11, bold=True, color=MUTE)
        _txt(slide.shapes, Inches(0.75), y + Inches(0.48), Inches(3.6), Inches(0.6),
             val, size=26, bold=True, color=INK)
        _txt(slide.shapes, Inches(4.3), y + Inches(0.55), Inches(1.7), Inches(0.4),
             sub, size=11, color=MUTE, align=PP_ALIGN.RIGHT)
    # top suppliers table + chart
    top = _supplier_totals(AP, GRN)[:8]
    if top:
        _txt(slide.shapes, Inches(6.4), Inches(1.2), Inches(6.4), Inches(0.4),
             "Top suppliers by actual spending (AED)", size=13, bold=True, color=INK)
        _chart(slide, XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.3), Inches(1.55),
               Inches(6.6), Inches(3.1),
               [k[:24] for k, _ in top][::-1], [("Actual", [round(v / 1e6, 2) for _, v in top][::-1])],
               colors=[TEAL], legend=False, data_labels=True, number_format='#,##0.0')
        _txt(slide.shapes, Inches(6.4), Inches(4.7), Inches(6.4), Inches(0.3),
             "Values in AED millions.", size=9, italic=True, color=MUTE)
    # concentration insight
    tot = ap_tot + grn_tot
    if top and tot:
        c1 = 100 * top[0][1] / tot
        _rect(slide.shapes, Inches(0.5), Inches(4.3), Inches(5.6), Inches(2.4), SOFT, line=LINE)
        _txt(slide.shapes, Inches(0.75), Inches(4.5), Inches(5.2), Inches(0.4),
             "Supplier concentration", size=13, bold=True, color=INK)
        _txt(slide.shapes, Inches(0.75), Inches(4.95), Inches(5.2), Inches(0.9),
             f"{c1:,.1f}%", size=38, bold=True, color=BAD if c1 > 50 else OK)
        _txt(slide.shapes, Inches(0.75), Inches(5.85), Inches(5.2), Inches(0.75),
             f"of all actual spending is with {top[0][0]}." +
             (" Heavily concentrated." if c1 > 50 else " Well spread across the supplier base."),
             size=12, color=MUTE)
    _foot(slide, meta)


def _open_slide(prs, ctx, PO, PR, O, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "05", "Open Obligations & Commitments")
    po_tot = sum(_num(r.get("open_aed")) for r in PO)
    pr_tot = sum(_num(r.get("amount_aed")) for r in PR)
    for i, (lbl, val, sub, accent) in enumerate([
        ("Open Obligations — Purchase Orders", money(po_tot), f"{len(PO):,} PO lines · GRN-netted", TEAL),
        ("Open Commitments — Purchase Requisitions", money(pr_tot), f"{len(PR):,} PR lines · reserved", GOLD)]):
        x = Inches(0.5 + i * 6.45)
        _rect(slide.shapes, x, Inches(1.3), Inches(6.15), Inches(1.35), WHITE, line=LINE)
        _rect(slide.shapes, x, Inches(1.3), Inches(6.15), Emu(int(0.07 * EMU_IN)), accent)
        _txt(slide.shapes, x + Inches(0.25), Inches(1.5), Inches(5.7), Inches(0.4),
             lbl.upper(), size=11, bold=True, color=MUTE)
        _txt(slide.shapes, x + Inches(0.25), Inches(1.9), Inches(5.7), Inches(0.6),
             val, size=28, bold=True, color=INK)
        _txt(slide.shapes, x + Inches(0.25), Inches(2.5), Inches(5.7), Inches(0.3),
             sub, size=11, color=MUTE)
    # two side-by-side top-5 tables
    po_rows = [[str(r.get("po_number") or ""), (str(r.get("supplier_name") or ""))[:24],
                money(r.get("open_aed"))] for r in PO[:6]]
    pr_rows = [[str(r.get("pr_number") or ""), (str(r.get("description") or r.get("project_name") or ""))[:24],
                money(r.get("amount_aed"))] for r in PR[:6]]
    if po_rows:
        _txt(slide.shapes, Inches(0.5), Inches(2.95), Inches(6), Inches(0.35),
             "Largest open purchase orders", size=12, bold=True, color=INK)
        _add_table(slide, Inches(0.4), Inches(3.35), Inches(6.2),
                   ["PO #", "Vendor", "Open (AED)"], po_rows,
                   col_w=[1.4, 3, 1.6], aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT],
                   height=Inches(2.9))
    if pr_rows:
        _txt(slide.shapes, Inches(6.95), Inches(2.95), Inches(6), Inches(0.35),
             "Largest open requisitions", size=12, bold=True, color=INK)
        _add_table(slide, Inches(6.85), Inches(3.35), Inches(6.1),
                   ["PR #", "Description", "Reserved (AED)"], pr_rows,
                   col_w=[1.4, 3, 1.7], aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT],
                   height=Inches(2.9))
    _foot(slide, meta)


def _insights(ctx, O, SEC, PRS, AP, GRN, PO, PR):
    """Same observations the Briefing Book computes, as (text, good?) bullets."""
    out = []
    p = ctx.get("params") or {}
    ga = ctx.get("generated_at", "") or ""
    try:
        gy = int(p["period"][3:]) if p.get("period") else int(ga[:4])
        gm = int(p["period"][:2]) if p.get("period") else int(ga[5:7])
    except (ValueError, KeyError, IndexError):
        gy, gm = 0, 0
    yr = int(_num(p.get("year")))
    elapsed = 100 if (yr and yr < gy) else (round(100 * gm / 12) if yr == gy else 0)
    util = O.get("utilization_pct") if O else None
    budget = _num(O.get("budget")) if O else 0
    if util is not None:
        u = _num(util)
        if u > elapsed + 10:
            out.append((f"Pacing: with {elapsed}% of the year elapsed, committed utilization is {u:,.1f}% "
                        "— running AHEAD of the calendar; funds may run out before year-end.", False))
        elif u < elapsed - 25:
            out.append((f"Pacing: utilization is {u:,.1f}% against {elapsed}% of the year elapsed "
                        "— running BEHIND the calendar; review delays or reallocation.", False))
        else:
            out.append((f"Pacing: utilization of {u:,.1f}% is broadly in line with the {elapsed}% of the year elapsed.", True))
    over = int(_num(O.get("over_budget_lines"))) if O else 0
    if over > 0:
        tight = ""
        if PRS:
            tight = f" The tightest is {PRS[0].get('project_number','')} with {money(PRS[0].get('fund_available'))} AED remaining."
        out.append((f"Budget pressure: {over} budget line(s) have spent more than their approved budget.{tight}", False))
    else:
        out.append(("Budget pressure: every budget line still has funds available.", True))
    top = _supplier_totals(AP, GRN)
    tot = sum(v for _, v in top)
    if top and tot:
        c1 = 100 * top[0][1] / tot
        if c1 > 50:
            out.append((f"Supplier concentration: the top supplier ({top[0][0]}) holds {c1:,.1f}% of actual spending "
                        "— heavily concentrated; a delivery or dispute there is a programme risk.", False))
        else:
            out.append((f"Supplier concentration: the top supplier holds {c1:,.1f}% of actual spending "
                        "— spending is well spread across the supplier base.", True))
    po_tot = sum(_num(r.get("open_aed")) for r in PO)
    pr_tot = sum(_num(r.get("amount_aed")) for r in PR)
    out.append((f"Delivery pipeline: {money(po_tot + pr_tot)} AED is committed-but-not-yet-delivered "
                f"({money(po_tot)} open POs + {money(pr_tot)} open PRs); it converts to actuals as goods are received and invoiced.", None))
    uninv = sum(_num(r.get("uninvoiced_aed")) for r in GRN)
    heavy = budget and uninv > 0.02 * budget
    if uninv > 0:
        out.append((f"Uninvoiced deliveries: {money(uninv)} AED of goods have been received without a matching supplier invoice "
                    + ("— a sizeable amount worth chasing with suppliers." if heavy else "— expected in the coming payment cycles."),
                    False if heavy else None))
    return out


def _insights_slide(prs, ctx, bullets, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "06", "Observations & Insights")
    y = 1.25
    for text, good in bullets:
        dot = OK if good is True else (BAD if good is False else TEAL)
        _rect(slide.shapes, Inches(0.55), Inches(y + 0.08), Inches(0.16), Inches(0.16), dot)
        box = slide.shapes.add_textbox(Inches(0.95), Inches(y), Inches(11.9), Inches(0.85))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(13.5)
        run.font.name = "Public Sans"
        run.font.color.rgb = INK
        y += 0.92
    _foot(slide, meta)


def _methodology_slide(prs, ctx, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "", "How the figures are defined", brand=TEAL_D)
    defs = [
        ("Actual — AP", "approved supplier invoices charged directly to projects (not linked to a purchase order), in AED."),
        ("Actual — GRN", "the value of goods and services received against purchase orders, in AED."),
        ("Open Obligation (PO)", "purchase-order lines with funds reserved, less what has already been received; finally-closed orders excluded."),
        ("Open Commitment (PR)", "purchase-requisition lines with funds reserved that have not yet become purchase orders."),
        ("Fund Available", "Budget − Actuals − Open PO − Open PR."),
        ("Utilization %", "(Actuals + Open PO + Open PR) ÷ Budget."),
    ]
    y = 1.4
    for term, desc in defs:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.9), Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        r1 = para.add_run()
        r1.text = term + " — "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.name = "Public Sans"
        r1.font.color.rgb = TEAL_D
        r2 = para.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.name = "Public Sans"
        r2.font.color.rgb = INK
        y += 0.72
    _txt(slide.shapes, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.7),
         "The scope — including any year-to-date period — applies the same filters, in the same way, as the "
         "Budget Utilization page this deck was generated from, so every figure here matches the page. "
         "Source: the i-Finance budget-utilization reporting layer, reflecting Oracle Fusion data as of the generation time.",
         size=11, italic=True, color=MUTE)
    _foot(slide, meta)


# ══ ENCUMBRANCES PENDING APPROVAL deck ═══════════════════════════════════════
# Consumes ENC_PENDING_BOOK's sections (reporting/db/23): overview / budget_ctx /
# aging / approvers / sectors / oldest / pr_register / po_register / unmatched —
# funds-RESERVED, non-zero pending PR/PO lines only, on the same GL Budget
# Utilization filter scope.
def _pn_cover(prs, ctx, O, B):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide.shapes, 0, 0, SW, Inches(2.35), TEAL)
    _rect(slide.shapes, 0, Inches(2.35), SW, Emu(int(0.08 * EMU_IN)), GOLD)
    _rect(slide.shapes, 0, Inches(7.18), SW, Inches(0.32), TEAL_D)
    p = ctx.get("params") or {}
    _txt(slide.shapes, Inches(0.7), Inches(0.55), Inches(12), Inches(0.4),
         "i-Finance · Reporting Platform", size=14, bold=True, color=RGBColor(0xCF, 0xE3, 0xEA))
    _txt(slide.shapes, Inches(0.7), Inches(1.0), Inches(12), Inches(1.2),
         "Encumbrances Pending Approval", size=42, bold=True, color=WHITE)
    sub = f"Budget Year {p.get('year','')}"
    if p.get("period"):
        sub += f"  —  YTD through {p['period']}"
    _txt(slide.shapes, Inches(0.7), Inches(2.6), Inches(12), Inches(0.5),
         sub, size=20, bold=True, color=TEAL_D)
    _txt(slide.shapes, Inches(0.7), Inches(3.25), Inches(12), Inches(0.5),
         "Scope:  " + _scope_line(p), size=13, color=MUTE)
    docs = _num(O.get("docs"))
    tiles = [
        ("Pending Documents", num(docs)),
        ("Pending Amount", money(O.get("amt_total"))),
        ("In Open Encumbrance", money(O.get("enc_open_aed"))),
        ("Avg Days Pending", num(O.get("avg_days"), 1)),
        ("Waiting > 30 Days", num(O.get("over30_docs"))),
    ]
    n = len(tiles)
    gap = Inches(0.2)
    tw = (SW - Inches(1.4) - gap * (n - 1)) / n
    x = Inches(0.7)
    for i, (lbl, val) in enumerate(tiles):
        _rect(slide.shapes, x, Inches(4.15), tw, Inches(1.35), SOFT, line=LINE)
        _rect(slide.shapes, x, Inches(4.15), tw, Emu(int(0.06 * EMU_IN)), GOLD if i == 4 else TEAL)
        _txt(slide.shapes, x, Inches(4.32), tw, Inches(0.35),
             lbl.upper(), size=10, bold=True, color=MUTE, align=PP_ALIGN.CENTER)
        _txt(slide.shapes, x, Inches(4.7), tw, Inches(0.7),
             val, size=22, bold=True, color=TEAL_D, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
        x += tw + gap
    _txt(slide.shapes, Inches(0.7), Inches(6.1), Inches(9), Inches(0.7),
         "Prepared by Financial Planning and Budgeting\nFinance Department",
         size=13, bold=True, color=INK)
    _txt(slide.shapes, Inches(0.7), Inches(7.2), Inches(12), Inches(0.28),
         f"Generated {ctx.get('generated_at','')} GST (UTC+04:00), Asia/Dubai   ·   {ctx.get('report_code','')}",
         size=9, color=RGBColor(0xCF, 0xE3, 0xEA), anchor=MSO_ANCHOR.MIDDLE)


def _pn_overview_slide(prs, ctx, O, B, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "01", "Pending Approval Overview")
    if not O or _num(O.get("docs")) == 0:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No funds-reserved documents are pending approval in this scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    docs = _num(O.get("docs"))
    pr_docs = _num(O.get("pr_docs"))
    po_docs = _num(O.get("po_docs"))
    amt = _num(O.get("amt_total"))
    amt_pr = _num(O.get("amt_pr"))
    amt_po = _num(O.get("amt_po"))
    tiles = [
        ("Pending Documents (PR+PO)", num(docs), f"PR {num(pr_docs)} · PO {num(po_docs)}", TEAL),
        ("Pending Amount (AED)", money(amt), f"PR {money(amt_pr)} · PO {money(amt_po)}", RGBColor(0x6C, 0x4A, 0xB6)),
        ("In Open Encumbrance (AED)", money(O.get("enc_open_aed")), f"{num(O.get('lines'))} pending lines", GOLD),
        ("Approval Aging", num(O.get("avg_days"), 1) + " d avg",
         f"longest {num(O.get('max_days'))} d · {num(O.get('over30_docs'))} over 30 d",
         BAD if _num(O.get("over30_docs")) > 0 else OK),
    ]
    n = len(tiles)
    gap = Inches(0.25)
    tw = (SW - Inches(1.0) - gap * (n - 1)) / n
    x = Inches(0.5)
    for lbl, val, sub, accent in tiles:
        _rect(slide.shapes, x, Inches(1.3), tw, Inches(1.9), WHITE, line=LINE)
        _rect(slide.shapes, x, Inches(1.3), Emu(int(0.08 * EMU_IN)), Inches(1.9), accent)
        _txt(slide.shapes, x + Inches(0.2), Inches(1.5), tw - Inches(0.3), Inches(0.55),
             lbl.upper(), size=11, bold=True, color=MUTE)
        _txt(slide.shapes, x + Inches(0.2), Inches(2.05), tw - Inches(0.3), Inches(0.7),
             val, size=26, bold=True, color=INK)
        _txt(slide.shapes, x + Inches(0.2), Inches(2.75), tw - Inches(0.3), Inches(0.4),
             sub, size=10.5, color=MUTE)
        x += tw + gap
    # left: PR vs PO amount doughnut
    comp = [(k, v) for k, v in [("Requisitions (PR)", amt_pr), ("Purchase Orders (PO)", amt_po)] if v > 0]
    if comp:
        _txt(slide.shapes, Inches(0.5), Inches(3.5), Inches(6), Inches(0.4),
             "Pending amount — PR vs PO", size=13, bold=True, color=INK)
        _chart(slide, XL_CHART_TYPE.DOUGHNUT, Inches(0.4), Inches(3.9), Inches(6.0), Inches(3.0),
               [c[0] for c in comp], [("Pending", [round(c[1], 2) for c in comp])],
               colors=[TEAL, GOLD], legend=True)
    # right: the budget frame the pending demand lands on
    _rect(slide.shapes, Inches(6.9), Inches(3.65), Inches(5.95), Inches(3.15), SOFT, line=LINE)
    _txt(slide.shapes, Inches(7.15), Inches(3.8), Inches(5.5), Inches(0.4),
         "The budget frame this demand lands on", size=13, bold=True, color=INK)
    frows = [
        ("Approved budget", money(B.get("budget"))),
        ("Actuals to date", money(B.get("actual_total"))),
        ("Open encumbrance (PR+PO)", money(B.get("open_encumbrance"))),
        ("Fund available", money(B.get("fund_available"))),
    ]
    y = 4.35
    for lbl, val in frows:
        _txt(slide.shapes, Inches(7.15), Inches(y), Inches(3.7), Inches(0.45),
             lbl, size=13, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        neg = lbl.startswith("Fund") and _num(B.get("fund_available")) < 0
        _txt(slide.shapes, Inches(10.6), Inches(y), Inches(2.1), Inches(0.45),
             val, size=15, bold=True, color=BAD if neg else TEAL_D,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += 0.6
    _foot(slide, meta)


def _pn_aging_slide(prs, ctx, O, AG, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "02", "Approval Aging")
    order = {"0-7": 1, "8-15": 2, "16-30": 3, "31+": 4}
    rows = sorted(AG, key=lambda r: order.get(str(r.get("bucket")), 9))
    if not rows:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No pending documents in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    cats = [str(r.get("bucket")) + " days" for r in rows]
    docs = [int(_num(r.get("docs"))) for r in rows]
    amt = [round(_num(r.get("amount_aed")) / 1e6, 2) for r in rows]
    _txt(slide.shapes, Inches(0.5), Inches(1.15), Inches(6), Inches(0.4),
         "Documents by days pending", size=13, bold=True, color=INK)
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(1.55),
           Inches(6.3), Inches(3.2), cats, [("Documents", docs)],
           colors=[TEAL], legend=False, data_labels=True)
    _txt(slide.shapes, Inches(6.95), Inches(1.15), Inches(6), Inches(0.4),
         "Pending amount by age (AED millions)", size=13, bold=True, color=INK)
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.85), Inches(1.55),
           Inches(6.0), Inches(3.2), cats, [("Amount (AED M)", amt)],
           colors=[GOLD], legend=False, data_labels=True, number_format='#,##0.0')
    # table under the charts
    hdr = ["Age bucket", "Documents", "PR", "PO", "Amount (AED)"]
    trows = [[str(r.get("bucket")) + " days", num(r.get("docs")), num(r.get("pr_docs")),
              num(r.get("po_docs")), money(r.get("amount_aed"))] for r in rows]
    over = int(_num(O.get("over30_docs"))) if O else 0
    if over:
        trows.append(["— of which over 30 days", num(over), "", "", money(O.get("over30_amt"))])
    aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4
    _add_table(slide, Inches(0.4), Inches(4.95), Inches(12.5), hdr, trows,
               col_w=[3, 1.4, 1, 1, 1.8], aligns=aligns, height=Inches(1.9))
    _foot(slide, meta)


def _pn_sector_slide(prs, ctx, SECR, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "03", "Pending Value by Sector")
    rows = sorted(SECR, key=lambda r: -_num(r.get("amount_aed")))[:8]
    if not rows:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No pending value in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    cats = [str(r.get("sector") or "—")[:24] for r in rows]
    amt = [round(_num(r.get("amount_aed")) / 1e6, 2) for r in rows]
    over = [round(_num(r.get("over30_aed")) / 1e6, 2) for r in rows]
    _chart(slide, XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(1.2),
           Inches(12.5), Inches(3.0), cats,
           [("Pending (AED M)", amt), ("Of which > 30 days (AED M)", over)],
           colors=[TEAL, GOLD], legend=True)
    hdr = ["Sector", "Documents", "Lines", "Pending (AED)", "> 30 days (AED)"]
    trows = [[str(r.get("sector") or "—"), num(r.get("docs")), num(r.get("lines")),
              money(r.get("amount_aed")), money(r.get("over30_aed"))] for r in rows]
    aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4
    _add_table(slide, Inches(0.4), Inches(4.4), Inches(12.5), hdr, trows,
               col_w=[3.4, 1.2, 1.2, 1.8, 1.8], aligns=aligns, height=Inches(2.4))
    _foot(slide, meta)


def _pn_approvers_slide(prs, ctx, APV, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "04", "Approver Follow-up")
    rows = APV[:12]
    if not rows:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No pending approvers in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    _txt(slide.shapes, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         "Who the pending documents are waiting with — highest pending value first.",
         size=12, color=MUTE)
    hdr = ["Pending with", "Documents", "Lines", "Pending (AED)", "Max days", "Oldest submitted"]
    trows = []
    for r in rows:
        trows.append([str(r.get("pending_with") or "(Unassigned)")[:40], num(r.get("docs")),
                      num(r.get("lines")), money(r.get("amount_aed")), num(r.get("max_days")),
                      _d(r.get("oldest_submitted"))])
    aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4 + [PP_ALIGN.LEFT]
    _add_table(slide, Inches(0.4), Inches(1.7), Inches(12.55), hdr, trows,
               col_w=[3.6, 1.2, 1, 1.6, 1, 1.6], aligns=aligns, height=Inches(5.0))
    _foot(slide, meta)


def _pn_oldest_slide(prs, ctx, OLD, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "05", "Longest-Waiting Documents")
    rows = OLD[:12]
    if not rows:
        _txt(slide.shapes, Inches(0.7), Inches(2), Inches(12), Inches(1),
             "No pending documents in scope.", size=16, color=MUTE)
        _foot(slide, meta)
        return
    _txt(slide.shapes, Inches(0.5), Inches(1.15), Inches(12), Inches(0.4),
         "The documents that have been waiting the longest — escalation candidates.",
         size=12, color=MUTE)
    hdr = ["Src", "Document #", "Description / Vendor", "Pending with", "Submitted",
           "Days", "Amount (AED)", "Sector"]
    trows = []
    for r in rows:
        trows.append([str(r.get("source") or ""), str(r.get("doc_number") or "")[:16],
                      str(r.get("description") or "")[:30], str(r.get("pending_with") or "")[:24],
                      _d(r.get("submitted_date")), num(r.get("pending_days")),
                      money(r.get("doc_amt")), str(r.get("sector") or "")[:18]])
    aligns = [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT,
              PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT]
    _add_table(slide, Inches(0.4), Inches(1.7), Inches(12.55), hdr, trows,
               col_w=[0.7, 1.6, 2.9, 2.1, 1.2, 0.8, 1.5, 1.6], aligns=aligns, height=Inches(5.0))
    _foot(slide, meta)


def _pn_insights(ctx, O, B, SECR, APV, OLD, UM):
    """Pending-approval observations, as (text, good?) bullets."""
    out = []
    docs = _num(O.get("docs"))
    amt = _num(O.get("amt_total"))
    if docs == 0:
        return [("No funds-reserved documents are pending approval in this scope.", True)]
    out.append((f"Backlog: {num(docs)} document(s) — {num(O.get('pr_docs'))} PR + {num(O.get('po_docs'))} PO — "
                f"worth {money(amt)} AED are awaiting approval; average wait {num(O.get('avg_days'),1)} days, "
                f"longest {num(O.get('max_days'))} days.", None))
    over = int(_num(O.get("over30_docs")))
    if over > 0:
        out.append((f"Aging pressure: {num(over)} document(s) worth {money(O.get('over30_amt'))} AED have been "
                    "pending more than 30 days — escalate these with the approvers first.", False))
    else:
        out.append(("Aging pressure: no document has been pending for more than 30 days.", True))
    tot_apv = sum(_num(r.get("amount_aed")) for r in APV)
    if APV and tot_apv:
        top = APV[0]
        c1 = 100 * _num(top.get("amount_aed")) / tot_apv
        out.append((f"Approver concentration: {str(top.get('pending_with') or '(Unassigned)')} holds "
                    f"{money(top.get('amount_aed'))} AED ({c1:,.1f}%) across {num(top.get('docs'))} document(s), "
                    f"waiting up to {num(top.get('max_days'))} days."
                    + (" A single approver is a bottleneck." if c1 > 40 else ""),
                    False if c1 > 40 else None))
    tot_sec = sum(_num(r.get("amount_aed")) for r in SECR)
    if SECR and tot_sec:
        s = sorted(SECR, key=lambda r: -_num(r.get("amount_aed")))[0]
        c1 = 100 * _num(s.get("amount_aed")) / tot_sec
        out.append((f"Sector concentration: {str(s.get('sector') or '—')} carries {money(s.get('amount_aed'))} AED "
                    f"({c1:,.1f}%) of the pending value.", None))
    budget = _num(B.get("budget"))
    fund = _num(B.get("fund_available"))
    if amt and (budget or fund):
        base = fund if fund > 0 else budget
        if base:
            pct = 100 * amt / base
            ref = "the available fund" if fund > 0 else "the approved budget"
            out.append((f"Budget context: the pending demand is {pct:,.1f}% of {ref} "
                        f"({money(base)} AED) for this scope.", None))
    um = _num(UM[0].get("docs")) if UM else 0
    # UM section is line-grain rows of un-matched documents; count distinct docs
    if UM:
        um_docs = len({(str(r.get("source")), str(r.get("doc_number"))) for r in UM})
        if um_docs:
            out.append((f"Extract coverage: {num(um_docs)} pending document(s) in the Fusion snapshot are not yet "
                        "matched to the budget extract, so they are outside the budget-scoped sections above "
                        "(listed in full in the Briefing Book annex).", None))
    return out


def _pn_methodology_slide(prs, ctx, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _band(slide, "", "How the figures are defined", brand=TEAL_D)
    defs = [
        ("Pending document", "a PR or PO awaiting approval in Fusion in the daily snapshot; counted once regardless of how many lines it has."),
        ("Scope rule", "funds-RESERVED, non-zero lines only — not-reserved documents and zero-value lines are excluded (the GL page still shows them)."),
        ("Pending amount", "the AED value of the pending distribution lines (PR reserved / PO reserved, GRN-netted)."),
        ("In open encumbrance", "the portion of that value still carried as an open commitment/obligation against the budget."),
        ("Days pending", "calendar days since the document was submitted for approval, as of the snapshot time."),
        ("Pending with", "the approver the document is currently waiting on."),
    ]
    y = 1.4
    for term, desc in defs:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(11.9), Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        r1 = para.add_run()
        r1.text = term + " — "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.name = "Public Sans"
        r1.font.color.rgb = TEAL_D
        r2 = para.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.name = "Public Sans"
        r2.font.color.rgb = INK
        y += 0.78
    _txt(slide.shapes, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.7),
         "The scope — including any year-to-date period and business-unit filter — applies the same filters, in the "
         "same way, as the Budget Utilization page this deck was generated from. Source: the daily Fusion "
         "pending-approval snapshot joined to the i-Finance budget-utilization reporting layer.",
         size=11, italic=True, color=MUTE)
    _foot(slide, meta)


def _enc_pending_deck(prs, sections, ctx, meta):
    S = _sections_by_key(sections)
    O = (_data(S.get("overview")) or [{}])[0] if S.get("overview") else {}
    B = (_data(S.get("budget_ctx")) or [{}])[0] if S.get("budget_ctx") else {}
    AG = _data(S.get("aging"))
    APV = _data(S.get("approvers"))
    SECR = _data(S.get("sectors"))
    OLD = _data(S.get("oldest"))
    UM = _data(S.get("unmatched"))
    _pn_cover(prs, ctx, O, B)
    _scope_slide(prs, ctx, meta)                 # Part 1 = Scope (selected parameters)
    _pn_aging_slide(prs, ctx, O, AG, meta)
    _pn_sector_slide(prs, ctx, SECR, meta)
    _pn_approvers_slide(prs, ctx, APV, meta)
    _pn_oldest_slide(prs, ctx, OLD, meta)
    _insights_slide(prs, ctx, _pn_insights(ctx, O, B, SECR, APV, OLD, UM), meta)
    _pn_methodology_slide(prs, ctx, meta)


# ── Budget Utilization deck ───────────────────────────────────────────────────
def _budget_util_deck(prs, sections, ctx, meta):
    S = _sections_by_key(sections)
    O = (_data(S.get("overview")) or [{}])[0] if S.get("overview") else {}
    SEC = _data(S.get("by_sector"))
    PRS = _data(S.get("pressure"))
    AP = _data(S.get("ap_lines"))
    GRN = _data(S.get("grn_lines"))
    PO = _data(S.get("open_po"))
    PR = _data(S.get("open_pr"))
    _cover(prs, ctx, O)
    _scope_slide(prs, ctx, meta)                 # Part 1 = Scope (selected parameters)
    _sector_slide(prs, ctx, SEC, meta)
    _pressure_slide(prs, ctx, PRS, meta)
    _actuals_slide(prs, ctx, AP, GRN, O, meta)
    _open_slide(prs, ctx, PO, PR, O, meta)
    _insights_slide(prs, ctx, _insights(ctx, O, SEC, PRS, AP, GRN, PO, PR), meta)
    _methodology_slide(prs, ctx, meta)


# ── entry point ──────────────────────────────────────────────────────────────
def build_deck(sections, ctx):
    """Build an executive .pptx deck from a MULTI report's sections.

    Dispatches on the report code: ENC_PENDING_BOOK → the pending-approval deck;
    everything else falls back to the Budget Utilization deck (the two share the
    same design language and slide helpers).
    """
    # content-slide footer suppressed per review (no bottom meta / GST line);
    # the cover keeps its generation stamp. Empty meta => _foot draws nothing.
    meta = ""
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    code = (ctx.get("report_code") or "").upper()
    keys = {(s or {}).get("key") for s in (sections or [])}
    if code == "ENC_PENDING_BOOK" or ("approvers" in keys and "pr_register" in keys):
        _enc_pending_deck(prs, sections, ctx, meta)
    else:
        _budget_util_deck(prs, sections, ctx, meta)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
